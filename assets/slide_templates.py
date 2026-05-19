"""
Templates pré-fabricados de slides - Código testado e funcional
A IA escolhe o template e personaliza cores, mas não gera código do zero
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import tempfile
from io import BytesIO


# ============================================================
# FUNÇÃO PARA GERAR GRÁFICOS COMO IMAGENS
# ============================================================

def gerar_grafico_slide(dados_mensais, cores=None):
    """
    Gera um gráfico de tendências financeiras como imagem para inserir no slide

    Args:
        dados_mensais: Lista de dicionários com dados mensais
        cores: Dicionário de cores (opcional)

    Returns:
        BytesIO: Imagem do gráfico em memória
    """
    try:
        import matplotlib
        matplotlib.use('Agg')  # Backend não-interativo
        import matplotlib.pyplot as plt
        from matplotlib.figure import Figure
    except ImportError:
        # Matplotlib não disponível, retornar None
        return None

    if not dados_mensais:
        return None

    # Preparar dados
    meses = [f"{d.get('mes', 0)}/{d.get('ano', 0)}" for d in dados_mensais]
    receita = [d.get('receita_bruta', 0) for d in dados_mensais]
    custos = [d.get('custos', 0) for d in dados_mensais]
    despesas = [d.get('despesas', 0) for d in dados_mensais]
    lucro = [d.get('lucro_operacional', 0) for d in dados_mensais]

    # Criar figura
    fig = Figure(figsize=(10, 6))
    ax1 = fig.add_subplot(2, 1, 1)
    ax2 = fig.add_subplot(2, 1, 2)

    # Definir cores baseadas no template
    if cores:
        cor_receita = cores.get('texto', RGBColor(0, 128, 0))
        cor_custos = cores.get('texto_secundario', RGBColor(255, 0, 0))
        cor_despesas = RGBColor(255, 165, 0)
        cor_lucro = cores.get('destaque', RGBColor(0, 0, 255))
    else:
        cor_receita = 'green'
        cor_custos = 'red'
        cor_despesas = 'orange'
        cor_lucro = 'blue'

    # Gráfico 1: Receita vs Custos vs Despesas
    ax1.plot(meses, receita, marker='o', label='Receita', color=cor_receita if isinstance(cor_receita, str) else 'green', linewidth=2)
    ax1.plot(meses, custos, marker='s', label='Custos', color=cor_custos if isinstance(cor_custos, str) else 'red', linewidth=2)
    ax1.plot(meses, despesas, marker='^', label='Despesas', color=cor_despesas if isinstance(cor_despesas, str) else 'orange', linewidth=2)
    ax1.set_title('Receita, Custos e Despesas por Mês', fontsize=12, fontweight='bold')
    ax1.set_ylabel('Valor (R$)', fontsize=10)
    ax1.legend()
    ax1.grid(True, alpha=0.3)
    ax1.tick_params(axis='x', rotation=45)

    # Gráfico 2: Lucro Operacional
    ax2.bar(meses, lucro, color=cor_lucro if isinstance(cor_lucro, str) else 'blue', alpha=0.7)
    ax2.set_title('Lucro Operacional por Mês', fontsize=12, fontweight='bold')
    ax2.set_ylabel('Lucro (R$)', fontsize=10)
    ax2.axhline(y=0, color='black', linestyle='-', linewidth=0.5)
    ax2.grid(True, alpha=0.3)
    ax2.tick_params(axis='x', rotation=45)

    fig.tight_layout()

    # Salvar em memória
    img_buffer = BytesIO()
    fig.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close(fig)

    return img_buffer


def adicionar_grafico_ao_slide(slide, img_buffer, x=0.5, y=1.5, width=9, height=4):
    """
    Adiciona um gráfico (imagem) ao slide

    Args:
        slide: Objeto slide do python-pptx
        img_buffer: BytesIO com a imagem do gráfico
        x: Posição X em polegadas
        y: Posição Y em polegadas
        width: Largura em polegadas
        height: Altura em polegadas
    """
    if img_buffer is None:
        return

    # Salvar imagem temporariamente
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
        tmp.write(img_buffer.getvalue())
        tmp_path = tmp.name

    try:
        # Adicionar imagem ao slide
        slide.shapes.add_picture(tmp_path, Inches(x), Inches(y), Inches(width), Inches(height))
    finally:
        # Remover arquivo temporário
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


# ============================================================
# FUNÇÃO DE ANÁLISE DE CENÁRIO FINANCEIRO
# ============================================================

def analisar_cenario_financeiro(dados_mensais):
    """
    Analisa os dados financeiros e retorna o cenário e recomendações adequadas

    Cenários:
    - LUCRO_ALTO: Margem > 20%
    - LUCRO_POSITIVO: Margem entre 5% e 20%
    - LUCRO_BAIXO: Margem entre 0% e 5%
    - PREJUIZO: Margem < 0%
    - CUSTOS_ALTOS: Custos > 70% da receita
    - CRESCIMENTO_ALTO: Receita crescendo > 15% mês a mês
    - CRESCIMENTO_BAIXO: Receita estagnada ou caindo
    """
    if not dados_mensais:
        return "SEM_DADOS", [], "Sem dados suficientes para análise."

    # Calcular métricas
    total_receita = sum(d.get('receita_bruta', 0) for d in dados_mensais)
    total_custos = sum(d.get('custos', 0) + d.get('despesas', 0) + d.get('impostos', 0) for d in dados_mensais)
    total_lucro = sum(d.get('lucro_operacional', 0) for d in dados_mensais)
    margem = (total_lucro / total_receita * 100) if total_receita > 0 else 0
    ratio_custos = (total_custos / total_receita * 100) if total_receita > 0 else 0

    # Analisar tendência de crescimento (comparar últimos 3 meses com os anteriores)
    if len(dados_mensais) >= 6:
        receita_recente = sum(d.get('receita_bruta', 0) for d in dados_mensais[-3:])
        receita_anterior = sum(d.get('receita_bruta', 0) for d in dados_mensais[-6:-3])
        if receita_anterior > 0:
            crescimento = ((receita_recente - receita_anterior) / receita_anterior) * 100
        else:
            crescimento = 0
    else:
        crescimento = 0

    # Determinar cenário
    cenario = "NEUTRO"
    if margem < 0:
        cenario = "PREJUIZO"
    elif margem < 5:
        cenario = "LUCRO_BAIXO"
    elif margem >= 30:
        cenario = "LUCRO_ALTO"
    elif ratio_custos > 70:
        cenario = "CUSTOS_ALTOS"
    elif crescimento > 15:
        cenario = "CRESCIMENTO_ALTO"
    elif crescimento < -5:
        cenario = "CRESCIMENTO_BAIXO"
    else:
        cenario = "LUCRO_POSITIVO"

    # Recomendações baseadas no cenário
    recomendacoes_por_cenario = {
        "PREJUIZO": [
            "1. Análise emergencial de estrutura de custos",
            "2. Revisão imediata de despesas operacionais",
            "3. Renegociação de contratos e fornecedores",
            "4. Planejamento tributário para reduzir carga fiscal",
            "5. Consultoria para reestruturação financeira"
        ],
        "LUCRO_BAIXO": [
            "1. Otimização de margens de lucro por produto/serviço",
            "2. Revisão de precificação e custos variáveis",
            "3. Identificação de desperdícios operacionais",
            "4. Análise de eficiência tributária",
            "5. Planejamento de expansão controlada"
        ],
        "LUCRO_POSITIVO": [
            "1. Manutenção de controle rigoroso de custos",
            "2. Reinvestimento estratégico em crescimento",
            "3. Diversificação de fontes de receita",
            "4. Planejamento tributário preventivo",
            "5. Reserva de contingência para imprevistos"
        ],
        "LUCRO_ALTO": [
            "1. Estruturação de holding para otimização fiscal",
            "2. Investimento em expansão e novos mercados",
            "3. Planejamento sucessório e patrimonial",
            "4. Diversificação de portfólio de investimentos",
            "5. Estratégias de hedge cambial e de commodities"
        ],
        "CUSTOS_ALTOS": [
            "1. Auditoria detalhada de estrutura de custos",
            "2. Renegociação de contratos de fornecedores",
            "3. Automatização de processos para reduzir desperdícios",
            "4. Terceirização de atividades não essenciais",
            "5. Revisão de logística e cadeia de suprimentos"
        ],
        "CRESCIMENTO_ALTO": [
            "1. Estruturação para escalabilidade do negócio",
            "2. Contratação e treinamento de equipe",
            "3. Investimento em tecnologia e infraestrutura",
            "4. Planejamento de expansão geográfica",
            "5. Gestão de fluxo de caixa para sustentar crescimento"
        ],
        "CRESCIMENTO_BAIXO": [
            "1. Análise de mercado e concorrência",
            "2. Revisão de estratégia de marketing e vendas",
            "3. Inovação em produtos/serviços",
            "4. Fidelização de base de clientes",
            "5. Parcerias estratégicas para alavancagem"
        ],
        "NEUTRO": [
            "1. Manutenção de controle financeiro rigoroso",
            "2. Monitoramento constante de indicadores",
            "3. Planejamento tributário preventivo",
            "4. Revisão periódica de estrutura de custos",
            "5. Consultoria contínua para compliance fiscal"
        ],
        "SEM_DADOS": [
            "1. Coleta e organização de dados financeiros",
            "2. Implementação de sistema de controle",
            "3. Treinamento da equipe em gestão financeira",
            "4. Definição de indicadores de performance",
            "5. Consultoria para estruturação financeira"
        ]
    }

    # Conclusão baseada no cenário
    conclusoes_por_cenario = {
        "PREJUIZO": f"Atenção: A empresa apresenta prejuízo de {abs(margem):.1f}%. É necessário agir urgentemente para reverter o cenário negativo através de reestruturação financeira e controle rigoroso de custos.",
        "LUCRO_BAIXO": f"A margem de lucro está em {margem:.1f}%, abaixo do ideal. Recomenda-se focar na otimização de custos e revisão de preciação para melhorar a rentabilidade.",
        "LUCRO_POSITIVO": f"A empresa apresenta margem saudável de {margem:.1f}%. Mantenha o controle financeiro e considere estratégias de expansão controlada para sustentar o crescimento.",
        "LUCRO_ALTO": f"Excelente desempenho com margem de {margem:.1f}%. É o momento ideal para planejamento estratégico, expansão e otimização patrimonial.",
        "CUSTOS_ALTOS": f"Os custos representam {ratio_custos:.1f}% da receita, acima do recomendado. Uma revisão estrutural da cadeia de custos é essencial para melhorar a rentabilidade.",
        "CRESCIMENTO_ALTO": f"A empresa apresenta crescimento expressivo de {crescimento:.1f}% no período recente. Prepare a estrutura para escalar o negócio de forma sustentável.",
        "CRESCIMENTO_BAIXO": f"A receita está estagnada ou em queda ({crescimento:.1f}%). É necessário revisar a estratégia de mercado e buscar novas oportunidades de crescimento.",
        "NEUTRO": f"A empresa apresenta desempenho estável com margem de {margem:.1f}%. Continue o monitoramento constante e mantenha as práticas de controle financeiro.",
        "SEM_DADOS": "Dados insuficientes para análise completa. Recomenda-se implementar um sistema robusto de controle financeiro para embasar decisões estratégicas."
    }

    return cenario, recomendacoes_por_cenario.get(cenario, recomendacoes_por_cenario["NEUTRO"]), conclusoes_por_cenario.get(cenario, conclusoes_por_cenario["NEUTRO"])


# ============================================================
# PALETAS DE CORES PRÉ-DEFINIDAS (Sem necessidade de Ollama)
# ============================================================

# 1. AUDITAR CLÁSSICO - Branco e Dourado (padrão)
CORES_AUDITAR = {
    'primaria': RGBColor(30, 58, 138),      # Azul marinho
    'secundaria': RGBColor(198, 103, 0),    # Marrom/Laranja forte (melhor contraste)
    'fundo': RGBColor(255, 255, 255),       # BRANCO
    'texto': RGBColor(33, 37, 41),          # Cinza escuro
    'texto_secundario': RGBColor(108, 117, 125),
    'destaque': RGBColor(212, 175, 55),     # Dourado
    'accent': RGBColor(212, 175, 55),       # Dourado para barras
}

# 2. AUDITAR ESCURO - Baseado na imagem (azul escuro/cinza)
CORES_AUDITAR_ESCURO = {
    'primaria': RGBColor(45, 55, 72),        # Azul acinzentado escuro
    'secundaria': RGBColor(160, 174, 192),   # Cinza azulado claro
    'fundo': RGBColor(30, 35, 45),           # FUNDO ESCURO (como na imagem)
    'texto': RGBColor(226, 232, 240),        # Texto claro
    'texto_secundario': RGBColor(160, 174, 192),
    'destaque': RGBColor(99, 179, 237),      # Azul claro destaque
    'accent': RGBColor(99, 179, 237),        # Azul para barras
}

# 3. AUDITAR AZUL PROFISSIONAL
CORES_AUDITAR_AZUL = {
    'primaria': RGBColor(30, 64, 175),       # Azul profundo
    'secundaria': RGBColor(147, 197, 253),   # Azul claro
    'fundo': RGBColor(15, 23, 42),            # Azul muito escuro
    'texto': RGBColor(241, 245, 249),         # Quase branco
    'texto_secundario': RGBColor(148, 163, 184),
    'destaque': RGBColor(96, 165, 250),      # Azul céu
    'accent': RGBColor(59, 130, 246),         # Azul vibrante
}

# 4. AUDITAR CORPORATIVO CINZA
CORES_AUDITAR_CINZA = {
    'primaria': RGBColor(55, 65, 81),        # Cinza escuro
    'secundaria': RGBColor(209, 213, 219),   # Cinza claro
    'fundo': RGBColor(17, 24, 39),            # Quase preto
    'texto': RGBColor(243, 244, 246),         # Cinza muito claro
    'texto_secundario': RGBColor(156, 163, 175),
    'destaque': RGBColor(251, 191, 36),      # Âmbar/Dourado
    'accent': RGBColor(245, 158, 11),         # Laranja dourado
}

# 5. AUDITAR VERDE EMPRESARIAL
CORES_AUDITAR_VERDE = {
    'primaria': RGBColor(6, 78, 59),          # Verde escuro
    'secundaria': RGBColor(52, 211, 153),     # Verde menta
    'fundo': RGBColor(236, 253, 245),         # Verde muito claro
    'texto': RGBColor(6, 78, 59),             # Verde escuro
    'texto_secundario': RGBColor(16, 185, 129),
    'destaque': RGBColor(5, 150, 105),        # Verde médio
    'accent': RGBColor(16, 185, 129),         # Verde esmeralda
}

# 6. AUDITAR VINHO ELEGANTE
CORES_AUDITAR_VINHO = {
    'primaria': RGBColor(88, 28, 135),        # Roxo vinho
    'secundaria': RGBColor(232, 121, 249),    # Rosa claro
    'fundo': RGBColor(250, 245, 255),         # Lilás claro
    'texto': RGBColor(88, 28, 135),           # Roxo escuro
    'texto_secundario': RGBColor(168, 85, 247),
    'destaque': RGBColor(192, 132, 252),      # Lilás
    'accent': RGBColor(147, 51, 234),         # Roxo vibrante
}

# Dicionário com todas as opções disponíveis
OPCOES_CORES = {
    'auditar_classico': ('Auditar Clássico (Branco/Dourado)', CORES_AUDITAR),
    'auditar_escuro': ('Auditar Escuro (Azul Escuro/Cinza)', CORES_AUDITAR_ESCURO),
    'auditar_azul': ('Auditar Azul Profissional', CORES_AUDITAR_AZUL),
    'auditar_cinza': ('Auditar Corporativo Cinza', CORES_AUDITAR_CINZA),
    'auditar_verde': ('Auditar Verde Empresarial', CORES_AUDITAR_VERDE),
    'auditar_vinho': ('Auditar Vinho Elegante', CORES_AUDITAR_VINHO),
}


def adicionar_logo(slide, x=None, y=0.2, largura=0.8):
    """Adiciona logo da Auditar no canto superior direito ou texto fallback"""
    try:
        # Se x não especificado, colocar no canto direito (slide width - logo width - margem)
        from pptx.util import Inches
        if x is None:
            x = 10 - largura - 0.3  # 10 inches slide width - logo - margem
        
        # Tentar encontrar logo em vários locais
        possiveis_paths = [
            os.path.join(os.path.dirname(__file__), "logo_auditar.png"),
            os.path.join(os.path.dirname(__file__), "assets", "logo_auditar.png"),
            os.path.join(os.getcwd(), "logo_auditar.png"),
            os.path.join(os.getcwd(), "assets", "logo_auditar.png"),
            os.path.join(os.path.dirname(os.path.dirname(__file__)), "logo_auditar.png"),
        ]
        
        logo_path = None
        for path in possiveis_paths:
            if os.path.exists(path):
                logo_path = path
                break
        
        if logo_path:
            slide.shapes.add_picture(logo_path, Inches(x), Inches(y), width=Inches(largura))
        else:
            # Fallback texto
            txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(largura), Inches(0.4))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "AUDITAR"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(212, 175, 55)  # Dourado
            p.alignment = PP_ALIGN.RIGHT
    except:
        pass  # Ignora erro se não conseguir adicionar


def adicionar_header_footer(slide, cores, slide_num=None, total_slides=None):
    """Adiciona header e footer padrão Auditar em todos os slides"""
    # Header - usar cor 'primaria' ou 'header' ou 'accent' como fallback
    cor_header = cores.get('header', cores.get('primaria', cores.get('accent', RGBColor(212, 175, 55))))
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.15))
    header.fill.solid()
    header.fill.fore_color.rgb = cor_header
    header.line.fill.background()

    # Footer - usar cor 'secundaria' ou 'footer' ou 'accent' como fallback
    cor_footer = cores.get('footer', cores.get('secundaria', cores.get('accent', RGBColor(212, 175, 55))))
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.47), Inches(10), Inches(0.15))
    footer.fill.solid()
    footer.fill.fore_color.rgb = cor_footer
    footer.line.fill.background()
    
    # Número do slide (opcional)
    if slide_num and total_slides:
        num_box = slide.shapes.add_textbox(Inches(9), Inches(5.15), Inches(0.8), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num}/{total_slides}"
        p.font.size = Pt(10)
        p.font.color.rgb = cores.get('texto_secundario', RGBColor(108, 117, 125))


def template_corporativo_escuro(prs, dados_mensais, nome_empresa, responsavel, cores=None):
    """Template: Corporativo Escuro (Azul Marinho + Dourado + BRANCO)"""
    if cores is None:
        cores = CORES_AUDITAR.copy()  # Usar cores padrão Auditar
    
    # Slide 1: Capa
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['fundo']
    background.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = nome_empresa
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = cores['texto']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Responsável: {responsavel}"
    p.font.size = Pt(20)
    p.font.color.rgb = cores['texto_secundario']
    p.alignment = PP_ALIGN.CENTER
    
    # Header do topo: Texto AUDITAR à esquerda + Logo à direita
    # Texto completo "AUDITAR - CONTABILIDADE CONSULTIVA S/S"
    auditar_top = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(6), Inches(0.6))
    tf = auditar_top.text_frame
    p = tf.paragraphs[0]
    p.text = "AUDITAR - CONTABILIDADE CONSULTIVA S/S"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = cores.get('accent', RGBColor(212, 175, 55))
    
    # Logo ao lado direito do texto
    adicionar_logo(slide, x=9.0, y=0.15, largura=0.7)
    
    # Header/footer
    adicionar_header_footer(slide, cores, 1, 8)
    
    # Slide 2: Dados Mensais
    if dados_mensais:
        slide = prs.slides.add_slide(slide_layout)
        
        # Fundo
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = cores['fundo']
        background.line.fill.background()
        
        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Dados Mensais"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = cores['texto']
        
        # Cards arredondados para cada mês (como nos cenários)
        for idx, dado in enumerate(dados_mensais[:6]):  # Max 6 meses no slide 2
            row = idx // 2
            col = idx % 2
            x = Inches(0.4 + col * 4.8)
            y = Inches(1.3 + row * 1.15)

            # Card arredondado
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1.0))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(248, 249, 250)  # Cinza claro
            card.line.color.rgb = cores.get('accent', RGBColor(212, 175, 55))
            card.line.width = Pt(2)


            # Mês/Ano
            mes_ano = f"{dado.get('mes', '01')}/{dado.get('ano', '2024')}"
            mes_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(4.1), Inches(0.35))
            tf = mes_box.text_frame
            p = tf.paragraphs[0]
            p.text = mes_ano
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

            # Receita
            rec_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.55), Inches(2), Inches(0.3))
            tf = rec_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Rec: R$ {dado.get('receita_bruta', 0):,.0f}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

            # Lucro
            luc_box = slide.shapes.add_textbox(x + Inches(2.2), y + Inches(0.55), Inches(2), Inches(0.3))
            tf = luc_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Luc: R$ {dado.get('lucro_operacional', 0):,.0f}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

        # Logo + header/footer slide 2
        adicionar_logo(slide)
        adicionar_header_footer(slide, cores, 2, 8 if len(dados_mensais) > 6 else 7)

    # Slide 2b: Dados Mensais (parte 2) - apenas se houver mais de 6 meses
    if len(dados_mensais) > 6:
        slide = prs.slides.add_slide(slide_layout)
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = cores['fundo']
        background.line.fill.background()

        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Dados Mensais (Continuação)"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = cores['texto']

        # Cards arredondados para meses 7-12
        for idx, dado in enumerate(dados_mensais[6:12]):  # Meses 7-12
            row = idx // 2
            col = idx % 2
            x = Inches(0.4 + col * 4.8)
            y = Inches(1.3 + row * 1.15)

            # Card arredondado
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(1.0))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(248, 249, 250)  # Cinza claro
            card.line.color.rgb = cores.get('accent', RGBColor(212, 175, 55))
            card.line.width = Pt(2)

            # Mês/Ano
            mes_ano = f"{dado.get('mes', '01')}/{dado.get('ano', '2024')}"
            mes_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(4.1), Inches(0.35))
            tf = mes_box.text_frame
            p = tf.paragraphs[0]
            p.text = mes_ano
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

            # Receita
            rec_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.55), Inches(2), Inches(0.3))
            tf = rec_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Rec: R$ {dado.get('receita_bruta', 0):,.0f}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

            # Lucro
            luc_box = slide.shapes.add_textbox(x + Inches(2.3), y + Inches(0.55), Inches(2), Inches(0.3))
            tf = luc_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Luc: R$ {dado.get('lucro_operacional', 0):,.0f}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

        # Logo + header/footer slide 2b
        adicionar_logo(slide)
        adicionar_header_footer(slide, cores, 3, 8)

    # Slide 3: Resumo Financeiro
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['fundo']
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RESUMO FINANCEIRO"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = cores['destaque']
    
    # Calcular totais
    total_receita = sum(d.get('receita_bruta', 0) for d in dados_mensais) if dados_mensais else 0
    total_lucro = sum(d.get('lucro_operacional', 0) for d in dados_mensais) if dados_mensais else 0
    margem = (total_lucro / total_receita * 100) if total_receita else 0
    
    # Cards de métricas
    metrics = [
        ("RECEITA TOTAL", f"R$ {total_receita:,.2f}", cores['secundaria']),
        ("LUCRO TOTAL", f"R$ {total_lucro:,.2f}", cores['destaque']),
        ("MARGEM", f"{margem:.1f}%", cores['texto']),
    ]
    
    for idx, (label, value, color) in enumerate(metrics):
        y = Inches(1.5 + idx * 1.2)
        
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(4), Inches(0.5))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = cores['texto_secundario']
        
        # Valor
        val_box = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.4), Inches(4), Inches(0.6))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
    
    # Logo + header/footer slide 3
    adicionar_logo(slide)
    adicionar_header_footer(slide, cores, 3, 8)
    
    # Slide 4: Análise
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['fundo']
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ANÁLISE TRIBUTÁRIA"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = cores['destaque']
    
    textos = [
        "Revisão da carga tributária atual",
        "Análise de regimes tributários disponíveis",
        "Otimização de custos dedutíveis",
        "Planejamento de fluxo de caixa",
    ]
    
    y = Inches(1.5)
    for texto in textos:
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {texto}"
        p.font.size = Pt(18)
        p.font.color.rgb = cores['texto']
        y += Inches(0.8)
    
    # Logo + header/footer slide 4
    adicionar_logo(slide)
    adicionar_header_footer(slide, cores, 4, 8)
    
    # Slide 5: CENÁRIOS
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['fundo']
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CENÁRIOS PROJETADOS"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = cores['destaque']
    
    # Calcular projeções
    total_receita = sum(d.get('receita_bruta', 0) for d in dados_mensais) if dados_mensais else 0
    total_lucro = sum(d.get('lucro_operacional', 0) for d in dados_mensais) if dados_mensais else 0
    n_meses = len(dados_mensais) if dados_mensais else 1
    
    media_receita = total_receita / n_meses
    media_lucro = total_lucro / n_meses
    
    cenarios = [
        ("Cenário Otimista", media_receita * 1.15, media_lucro * 1.25, cores['destaque']),
        ("Cenário Realista", media_receita, media_lucro, cores['primaria']),
        ("Cenário Conservador", media_receita * 0.85, media_lucro * 0.75, cores['secundaria']),
    ]
    
    for idx, (nome, rec, luc, cor) in enumerate(cenarios):
        y = Inches(1.4 + idx * 1.3)
        
        # Card do cenário
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(9), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 249, 250)
        card.line.color.rgb = cor
        card.line.width = Pt(2)
        

        
        # Nome do cenário - USAR CORES APROPRIADAS PARA CONTRASTE
        nome_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.1), Inches(8.6), Inches(0.4))
        tf = nome_box.text_frame
        p = tf.paragraphs[0]
        p.text = nome
        p.font.size = Pt(16)
        p.font.bold = True
        # Usar cor que contrasta bem com fundo branco do card
        p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

        # Valores
        val_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.5), Inches(8.6), Inches(0.5))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Receita: R$ {rec:,.2f} | Lucro: R$ {luc:,.2f}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste
    
    # Logo + header/footer slide 5
    adicionar_logo(slide)
    adicionar_header_footer(slide, cores, 5, 9 if len(dados_mensais) > 6 else 8)

    # Slide 6: Gráfico de Tendências
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['fundo']
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GRÁFICO DE TENDÊNCIAS"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = cores['destaque']

    # Gerar e adicionar gráfico
    grafico_img = gerar_grafico_slide(dados_mensais, cores)
    if grafico_img:
        adicionar_grafico_ao_slide(slide, grafico_img, x=0.5, y=1.2, width=8, height=3.5)
    else:
        # Matplotlib não disponível, mostrar mensagem
        msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Gráfico não disponível (matplotlib não instalado)"
        p.font.size = Pt(14)
        p.font.color.rgb = cores['texto']
        p.alignment = PP_ALIGN.CENTER

    # Logo + header/footer slide 6
    adicionar_logo(slide)
    adicionar_header_footer(slide, cores, 6, 9 if len(dados_mensais) > 6 else 8)

    # Slide 7: Recomendações (baseadas no cenário financeiro)
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['fundo']
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RECOMENDAÇÕES"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = cores['destaque']

    # Analisar cenário e obter recomendações dinâmicas
    cenario, recs, _ = analisar_cenario_financeiro(dados_mensais)

    y = Inches(1.5)
    for rec in recs:
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = rec
        p.font.size = Pt(16)
        p.font.color.rgb = cores['texto']
        y += Inches(0.7)

    # Logo + header/footer slide 7
    adicionar_logo(slide)
    adicionar_header_footer(slide, cores, 7, 9 if len(dados_mensais) > 6 else 8)

    # Slide 8: Conclusão (baseada no cenário financeiro)
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['fundo']
    background.line.fill.background()

    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = cores['secundaria']
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CONCLUSÃO"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = cores['destaque']
    p.alignment = PP_ALIGN.CENTER

    # Conclusão dinâmica baseada no cenário
    _, _, conclusao = analisar_cenario_financeiro(dados_mensais)

    concl_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
    tf = concl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = conclusao
    p.font.size = Pt(16)
    p.font.color.rgb = cores['texto']
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = nome_empresa.upper()
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = cores['texto']
    p.alignment = PP_ALIGN.CENTER

    resp_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.5))
    tf = resp_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Responsável: {responsavel}"
    p.font.size = Pt(14)
    p.font.color.rgb = cores['texto_secundario']
    p.alignment = PP_ALIGN.CENTER
    
    # Barra inferior
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.47), prs.slide_width, Inches(0.15))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = cores['secundaria']
    bar2.line.fill.background()



    # Logo + header/footer slide 8
    adicionar_logo(slide)
    adicionar_header_footer(slide, cores, 8, 9 if len(dados_mensais) > 6 else 8)
    
    return prs


def template_minimalista_branco(prs, dados_mensais, nome_empresa, responsavel, cores=None):
    """Template: Minimalista Branco com cores de destaque"""
    if cores is None:
        cores = {
            'primaria': RGBColor(51, 51, 51),       # Cinza escuro
            'secundaria': RGBColor(220, 53, 69),    # Vermelho
            'fundo': RGBColor(255, 255, 255),       # Branco
            'texto': RGBColor(33, 37, 41),          # Cinza escuro
            'texto_secundario': RGBColor(108, 117, 125),
            'destaque': RGBColor(40, 167, 69)      # Verde
        }
    
    # Slide 1: Capa
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fundo branco
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['fundo']
    background.line.fill.background()
    
    # Linha decorativa superior
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    line.fill.solid()
    line.fill.fore_color.rgb = cores['secundaria']
    line.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = nome_empresa
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = cores['texto']
    p.alignment = PP_ALIGN.CENTER
    
    # Linha decorativa abaixo do título
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.2), Inches(6), Inches(0.05))
    line2.fill.solid()
    line2.fill.fore_color.rgb = cores['destaque']
    line2.line.fill.background()
    
    # Responsável
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Responsável: {responsavel}"
    p.font.size = Pt(18)
    p.font.color.rgb = cores['texto_secundario']
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Dados
    if dados_mensais:
        slide = prs.slides.add_slide(slide_layout)
        
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = cores['fundo']
        background.line.fill.background()
        
        # Cabeçalho com cor
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = cores['secundaria']
        header.line.fill.background()
        
        # Título sobreposto
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Análise Mensal"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Cards com dados
        for idx, dado in enumerate(dados_mensais[:4]):
            row = idx // 2
            col = idx % 2
            x = Inches(0.5 + col * 4.7)
            y = Inches(1.5 + row * 2.5)
            
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.5), Inches(2.2))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(248, 249, 250)
            card.line.color.rgb = cores['destaque']
            
            # Mês/Ano
            mes_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(4.1), Inches(0.5))
            tf = mes_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{dado.get('mes', '01')}/{dado.get('ano', '2024')}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

            # Receita
            rec_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), Inches(4.1), Inches(0.4))
            tf = rec_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Receita: R$ {dado.get('receita_bruta', 0):,.2f}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste

            # Lucro
            luc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.3), Inches(4.1), Inches(0.4))
            tf = luc_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Lucro: R$ {dado.get('lucro_operacional', 0):,.2f}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(33, 37, 41)  # Cinza escuro para contraste
    
    # Slides 3-7: Resumo, Análise, Gráfico, Recomendações, Conclusão (simplificados)
    total_receita = sum(d.get('receita_bruta', 0) for d in dados_mensais) if dados_mensais else 0
    total_lucro = sum(d.get('lucro_operacional', 0) for d in dados_mensais) if dados_mensais else 0

    # Analisar cenário financeiro para recomendações e conclusão
    cenario, recs_dinamicas, conclusao_dinamica = analisar_cenario_financeiro(dados_mensais)

    for slide_num, titulo in enumerate(["RESUMO", "ANÁLISE", "GRÁFICO", "RECOMENDAÇÕES", "CONCLUSÃO"], 3):
        slide = prs.slides.add_slide(slide_layout)

        # Fundo
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = cores['fundo']
        background.line.fill.background()

        # Cabeçalho
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        header.fill.solid()
        header.fill.fore_color.rgb = cores['secundaria'] if slide_num % 2 == 0 else cores['destaque']
        header.line.fill.background()

        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        if slide_num == 3:  # Resumo
            # Métricas
            margem = (total_lucro / total_receita * 100) if total_receita else 0
            metrics = [
                ("Receita", f"R$ {total_receita:,.2f}"),
                ("Lucro", f"R$ {total_lucro:,.2f}"),
                ("Margem", f"{margem:.1f}%"),
            ]
            for i, (label, value) in enumerate(metrics):
                y = Inches(1.8 + i * 1)
                lbl_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(4), Inches(0.5))
                tf = lbl_box.text_frame
                p = tf.paragraphs[0]
                p.text = label
                p.font.size = Pt(14)
                p.font.color.rgb = cores['texto_secundario']

                val_box = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.4), Inches(4), Inches(0.6))
                tf = val_box.text_frame
                p = tf.paragraphs[0]
                p.text = value
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = cores['secundaria'] if i == 0 else cores['destaque'] if i == 1 else cores['texto']

        elif slide_num == 4:  # Gráfico
            # Gerar e adicionar gráfico
            grafico_img = gerar_grafico_slide(dados_mensais, cores)
            if grafico_img:
                adicionar_grafico_ao_slide(slide, grafico_img, x=0.5, y=1.2, width=8, height=3.5)
            else:
                # Matplotlib não disponível, mostrar mensagem
                msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
                tf = msg_box.text_frame
                p = tf.paragraphs[0]
                p.text = "Gráfico não disponível (matplotlib não instalado)"
                p.font.size = Pt(14)
                p.font.color.rgb = cores['texto']
                p.alignment = PP_ALIGN.CENTER

        elif slide_num == 6:  # Recomendações
            # Recomendações dinâmicas baseadas no cenário
            y = Inches(1.5)
            for rec in recs_dinamicas:
                box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.5))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = rec
                p.font.size = Pt(16)
                p.font.color.rgb = cores['texto']
                y += Inches(0.7)

        elif slide_num == 7:  # Conclusão
            # Conclusão dinâmica
            concl_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
            tf = concl_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = conclusao_dinamica
            p.font.size = Pt(16)
            p.font.color.rgb = cores['texto']
            p.alignment = PP_ALIGN.CENTER

            # Título central
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = nome_empresa
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = cores['texto']
            p.alignment = PP_ALIGN.CENTER

            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(0.6))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Responsável: {responsavel}"
            p.font.size = Pt(14)
            p.font.color.rgb = cores['texto_secundario']
            p.alignment = PP_ALIGN.CENTER
    
    return prs


def template_moderno_gradiente(prs, dados_mensais, nome_empresa, responsavel, cores=None):
    """Template: Moderno com design clean e elementos visuais"""
    if cores is None:
        cores = {
            'primaria': RGBColor(67, 97, 238),      # Azul vibrante
            'secundaria': RGBColor(114, 9, 183),    # Roxo
            'fundo': RGBColor(248, 249, 250),      # Cinza claro
            'texto': RGBColor(33, 37, 41),
            'texto_secundario': RGBColor(108, 117, 125),
            'destaque': RGBColor(6, 214, 160)       # Verde menta
        }
    
    # Slide 1: Capa moderna
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fundo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = cores['primaria']
    background.line.fill.background()
    
    # Elemento decorativo - círculo
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-2), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = cores['secundaria']
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = nome_empresa
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Linha decorativa
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4), Inches(2), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = cores['destaque']
    line.line.fill.background()
    
    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(8), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Análise Financeira | {responsavel}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Slide 2: Resumo
    if dados_mensais:
        slide = prs.slides.add_slide(slide_layout)
        
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = cores['fundo']
        background.line.fill.background()
        
        # Header colorido
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
        header.fill.solid()
        header.fill.fore_color.rgb = cores['primaria']
        header.line.fill.background()
        
        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Resumo Financeiro"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Cards de resumo
        total_receita = sum(d.get('receita_bruta', 0) for d in dados_mensais)
        total_lucro = sum(d.get('lucro_operacional', 0) for d in dados_mensais)
        margem = (total_lucro / total_receita * 100) if total_receita else 0
        
        metrics = [
            ("Receita Total", f"R$ {total_receita:,.2f}", cores['primaria']),
            ("Lucro Total", f"R$ {total_lucro:,.2f}", cores['destaque']),
            ("Margem Média", f"{margem:.1f}%", cores['secundaria']),
        ]
        
        for idx, (label, value, color) in enumerate(metrics):
            x = Inches(0.5 + idx * 3.2)
            
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), Inches(3), Inches(2))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = color
            
            # Barra colorida no topo
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2), Inches(3), Inches(0.15))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            
            # Label
            label_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.4), Inches(2.6), Inches(0.5))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(14)
            p.font.color.rgb = cores['texto_secundario']
            
            # Valor
            val_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.9), Inches(2.6), Inches(0.8))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = color
    
    # Slides 3-6 adicionais
    # Analisar cenário financeiro para recomendações e conclusão
    cenario, recs_dinamicas, conclusao_dinamica = analisar_cenario_financeiro(dados_mensais)

    for slide_num, titulo in enumerate(["ANÁLISE", "DADOS MENSAIS", "GRÁFICO", "RECOMENDAÇÕES", "CONCLUSÃO"], 3):
        slide = prs.slides.add_slide(slide_layout)

        # Fundo
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = cores['fundo']
        background.line.fill.background()

        # Header colorido
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = cores['primaria'] if slide_num % 2 else cores['secundaria']
        header.line.fill.background()

        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        if slide_num == 4:  # Gráfico
            # Gerar e adicionar gráfico
            grafico_img = gerar_grafico_slide(dados_mensais, cores)
            if grafico_img:
                adicionar_grafico_ao_slide(slide, grafico_img, x=0.5, y=1.2, width=8, height=3.5)
            else:
                # Matplotlib não disponível, mostrar mensagem
                msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
                tf = msg_box.text_frame
                p = tf.paragraphs[0]
                p.text = "Gráfico não disponível (matplotlib não instalado)"
                p.font.size = Pt(14)
                p.font.color.rgb = cores['texto']
                p.alignment = PP_ALIGN.CENTER

        elif slide_num == 6:  # Recomendações
            # Recomendações dinâmicas baseadas no cenário
            y = Inches(1.5)
            for rec in recs_dinamicas:
                box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.5))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = rec
                p.font.size = Pt(16)
                p.font.color.rgb = cores['texto']
                y += Inches(0.7)

        elif slide_num == 7:  # Conclusão
            # Conclusão dinâmica
            concl_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
            tf = concl_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = conclusao_dinamica
            p.font.size = Pt(16)
            p.font.color.rgb = cores['texto']
            p.alignment = PP_ALIGN.CENTER

            # Centralizado
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = nome_empresa
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = cores['primaria']
            p.alignment = PP_ALIGN.CENTER

            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.6))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Responsável: {responsavel}"
            p.font.size = Pt(14)
            p.font.color.rgb = cores['texto_secundario']
            p.alignment = PP_ALIGN.CENTER
    
    return prs


# Mapeamento de templates
TEMPLATES = {
    "corporativo_escuro": template_corporativo_escuro,
    "minimalista_branco": template_minimalista_branco,
    "moderno_gradiente": template_moderno_gradiente,
}


def aplicar_template(nome_template, prs, dados_mensais, nome_empresa, responsavel, cores_personalizadas=None):
    """Aplica um template pré-fabricado com cores personalizadas"""
    template_func = TEMPLATES.get(nome_template, template_corporativo_escuro)
    return template_func(prs, dados_mensais, nome_empresa, responsavel, cores_personalizadas)
