"""
GERADOR DE APRESENTAÇÕES V2.0 - ESTILO PDF
Cria slides PowerPoint parecidos com o design do PDF de exemplo
"""

import os
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================
# PALETA DE CORES - Estilo PDF com Dourado e Branco (Primárias)
# ============================================================

# Dicionário de cores atual (pode ser sobrescrito)
CORES = {
    'COR_PRETO': RGBColor(0x2D, 0x2D, 0x2D),       # Textos principais
    'COR_DOURADO': RGBColor(0xB8, 0x8F, 0x00),      # Dourado corporativo
    'COR_BRANCO': RGBColor(0xFF, 0xFF, 0xFF),       # Fundo principal
    'COR_CINZA_ESC': RGBColor(0x4A, 0x4A, 0x4A),    # Headers estilo PDF
    'COR_CINZA_MED': RGBColor(0x7A, 0x7A, 0x7A),    # Textos secundários
    'COR_TEXTO_MUTED': RGBColor(0xA0, 0xA0, 0xA0),  # Textos secundários
}

# Atalhos para compatibilidade
COR_PRETO = CORES['COR_PRETO']
COR_DOURADO = CORES['COR_DOURADO']
COR_BRANCO = CORES['COR_BRANCO']
COR_CINZA_ESC = CORES['COR_CINZA_ESC']
COR_CINZA_MED = CORES['COR_CINZA_MED']
COR_TEXTO_MUTED = CORES['COR_TEXTO_MUTED']

def atualizar_cores(cores_personalizadas):
    """Atualiza as cores globais com base nas cores personalizadas"""
    global COR_PRETO, COR_DOURADO, COR_BRANCO, COR_CINZA_ESC, COR_CINZA_MED, COR_TEXTO_MUTED
    
    if cores_personalizadas:
        # Mapear cores personalizadas para as constantes
        CORES['COR_PRETO'] = cores_personalizadas.get('texto', RGBColor(0x2D, 0x2D, 0x2D))
        CORES['COR_DOURADO'] = cores_personalizadas.get('accent', RGBColor(0xB8, 0x8F, 0x00))
        CORES['COR_BRANCO'] = cores_personalizadas.get('fundo', RGBColor(0xFF, 0xFF, 0xFF))
        CORES['COR_CINZA_ESC'] = cores_personalizadas.get('primaria', RGBColor(0x4A, 0x4A, 0x4A))
        CORES['COR_CINZA_MED'] = cores_personalizadas.get('texto_secundario', RGBColor(0x7A, 0x7A, 0x7A))
        CORES['COR_TEXTO_MUTED'] = cores_personalizadas.get('texto_secundario', RGBColor(0xA0, 0xA0, 0xA0))
        
        # Atualizar atalhos
        COR_PRETO = CORES['COR_PRETO']
        COR_DOURADO = CORES['COR_DOURADO']
        COR_BRANCO = CORES['COR_BRANCO']
        COR_CINZA_ESC = CORES['COR_CINZA_ESC']
        COR_CINZA_MED = CORES['COR_CINZA_MED']
        COR_TEXTO_MUTED = CORES['COR_TEXTO_MUTED']
        
        print(f"🎨 Cores atualizadas: Fundo={CORES['COR_BRANCO']}, Destaque={CORES['COR_DOURADO']}")

def formatar_moeda(valor):
    return f"R$ {valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

def adicionar_fundo(slide, cor=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = cor or COR_BRANCO

def adicionar_retangulo(slide, x, y, w, h, cor_fill, cor_linha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor_fill
    if cor_linha:
        shape.line.color.rgb = cor_linha
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

from pptx.enum.text import MSO_ANCHOR

def adicionar_texto(slide, texto, x, y, w, h, tamanho=14, cor=None, bold=False,
                    italic=False, alinhamento=PP_ALIGN.LEFT, fonte="Calibri", vertical_center=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    if vertical_center:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = alinhamento
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.color.rgb = cor or COR_PRETO
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = fonte
    return txBox

def adicionar_logo_auditar(slide, x=0.3, y=4.6, largura=0.8):
    """Adiciona logo da Auditar acima do footer dourado"""
    try:
        logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo_auditar.png")

        if os.path.exists(logo_path):
            # Logo posicionada sobre o footer dourado
            slide.shapes.add_picture(logo_path, Inches(x), Inches(y), width=Inches(largura))
        else:
            # Fallback para texto estilizado
            txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(largura), Inches(0.4))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "🏢 AUDITAR"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COR_DOURADO
            p.alignment = PP_ALIGN.LEFT
    except:
        pass

# ============================================================
# SLIDE 1 - CAPA ESTILO PDF
# ============================================================
def slide_capa_pdf(prs, nome_empresa, responsavel, data_emissao):
    """Slide no estilo do PDF - capa profissional"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide, COR_BRANCO)
    
    # Header dourado
    adicionar_retangulo(slide, 0, 0, 10, 0.15, COR_DOURADO)
    
    # Título principal fora do header
    adicionar_texto(slide, "PLANEJAMENTO TRIBUTÁRIO", 0.3, 0.5, 9.4, 0.8,
                    tamanho=36, cor=COR_PRETO, bold=True, fonte="Arial")
    
    # Subtítulo
    adicionar_texto(slide, "ANÁLISE ESTRATÉGICA E PROJEÇÕES", 0.3, 1.3, 9.4, 0.5,
                    tamanho=20, cor=COR_PRETO, bold=True, fonte="Arial")
    
    # Linha separadora
    adicionar_retangulo(slide, 0.3, 1.8, 9.4, 0.02, COR_DOURADO)
    
    # Nome da empresa (nome original)
    adicionar_texto(slide, nome_empresa.upper(), 0.3, 2.0, 9.4, 0.6,
                    tamanho=24, cor=COR_PRETO, bold=True, fonte="Arial", alinhamento=PP_ALIGN.CENTER)
    
    # Informações
    adicionar_texto(slide, f"Responsável: {responsavel}", 0.3, 2.7, 9.4, 0.4,
                    tamanho=12, cor=COR_TEXTO_MUTED, fonte="Arial")
    adicionar_texto(slide, f"Data: {data_emissao}", 0.3, 3.1, 9.4, 0.4,
                    tamanho=12, cor=COR_TEXTO_MUTED, fonte="Arial")
    
    # Logo sobre o footer
    adicionar_logo_auditar(slide)
    
    # Texto do footer (posicionado acima da barra)
    adicionar_texto(slide, "🏢 AUDITAR - CONTABILIDADE CONSULTIVA 🏢", 0.3, 4.85, 9.4, 0.3,
                    tamanho=10, cor=COR_DOURADO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Footer dourado
    adicionar_retangulo(slide, 0, 5.2, 10, 0.1, COR_DOURADO)

# ============================================================
# SLIDE 2 - SUMÁRIO EXECUTIVO ESTILO PDF
# ============================================================
def slide_sumario_pdf(prs, dados_mensais, nome_empresa):
    """Slide de sumário no estilo PDF"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide, COR_BRANCO)
    
    # Header dourado
    adicionar_retangulo(slide, 0, 0, 10, 0.12, COR_DOURADO)
    adicionar_texto(slide, "02", 0.4, 0.02, 1, 0.6, tamanho=28, cor=COR_BRANCO, bold=True, fonte="Arial")
    adicionar_texto(slide, "📊 SUMÁRIO EXECUTIVO", 1.3, 0.02, 8, 0.6,
                    tamanho=24, cor=COR_PRETO, bold=True, fonte="Arial")
    
    # Área de conteúdo
    adicionar_retangulo(slide, 0.3, 0.8, 9.4, 3.8, COR_BRANCO, COR_CINZA_MED)
    
    # Indicadores principais
    receita_t = sum(d["receita_bruta"] for d in dados_mensais)
    lucro_t = sum(d["lucro_operacional"] for d in dados_mensais)
    margem = (lucro_t / receita_t * 100) if receita_t > 0 else 0
    
    # Cards de indicadores
    indicadores = [
        ("Receita Total", formatar_moeda(receita_t), COR_DOURADO),
        ("Lucro Operacional", formatar_moeda(lucro_t), COR_DOURADO),
        ("Margem Líquida", f"{margem:.1f}%", COR_DOURADO)
    ]
    
    for i, (titulo, valor, cor) in enumerate(indicadores):
        x = 0.5 + i * 3.2
        y = 1.2
        
        # Card
        adicionar_retangulo(slide, x, y, 2.8, 1.2, COR_BRANCO, cor)
        adicionar_texto(slide, titulo, x + 0.1, y + 0.1, 2.6, 0.3,
                        tamanho=10, cor=cor, bold=True, fonte="Arial")
        adicionar_texto(slide, valor, x + 0.1, y + 0.4, 2.6, 0.6,
                        tamanho=16, cor=COR_PRETO, bold=True, fonte="Arial")
    
    # Análise rápida
    adicionar_texto(slide, "ANÁLISE RÁPIDA", 0.5, 2.6, 8.4, 0.3,
                    tamanho=12, cor=COR_PRETO, bold=True, fonte="Arial")
    
    texto_analise = f"""
A empresa {nome_empresa} apresenta performance financeira sólida 
com receita total de {formatar_moeda(receita_t)} e margem operacional de {margem:.1f}%.

Recomendações principais:
• Focar em otimização tributária
• Manter controle de custos
• Buscar oportunidades de crescimento
"""
    
    adicionar_texto(slide, texto_analise.strip(), 0.5, 2.9, 8.4, 1.5,
                    tamanho=10, cor=COR_CINZA_ESC, fonte="Arial")
    
    # Logo sobre o footer
    adicionar_logo_auditar(slide)
    
    # Texto do footer (posicionado acima da barra)
    adicionar_texto(slide, "🏢 AUDITAR - CONTABILIDADE CONSULTIVA 🏢", 0.3, 4.85, 9.4, 0.3,
                    tamanho=10, cor=COR_DOURADO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Footer dourado
    adicionar_retangulo(slide, 0, 5.2, 10, 0.1, COR_DOURADO)

# ============================================================
# SLIDE 3 - CENÁRIOS ESTILO PDF
# ============================================================
def slide_cenarios_pdf(prs, dados_mensais, cenarios):
    """Slide de cenários no estilo PDF"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide, COR_BRANCO)
    
    # Header dourado
    adicionar_retangulo(slide, 0, 0, 10, 0.12, COR_DOURADO)
    adicionar_texto(slide, "03", 0.4, 0.02, 1, 0.6, tamanho=28, cor=COR_BRANCO, bold=True, fonte="Arial")
    
    # Título fora do header
    adicionar_texto(slide, "📈 ANÁLISE DE CENÁRIOS", 1.3, 0.02, 8, 0.6,
                    tamanho=24, cor=COR_PRETO, bold=True, fonte="Arial")
    
    # Tabela de cenários estilo profissional
    cabecalhos = ["Cenário", "Variação", "Receita", "Lucro", "Imposto", "Margem"]
    larguras = [1.5, 1.2, 1.8, 1.8, 1.5, 1.2]
    cores = [COR_DOURADO, COR_DOURADO, COR_DOURADO]
    
    # Header da tabela
    adicionar_retangulo(slide, 0.3, 0.8, 9.4, 0.4, COR_CINZA_ESC)
    
    for i, (cab, larg) in enumerate(zip(cabecalhos, larguras)):
        x = 0.3 + sum(larguras[:i])
        adicionar_texto(slide, cab, x + 0.05, 0.85, larg, 0.3,
                        tamanho=10, cor=COR_BRANCO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
    
    # Dados dos cenários
    nomes_cenarios = ["Otimista", "Realista", "Pessimista"]
    variacoes = ["+30%", "+10%", "-20%"]
    
    for i, (cenario, variacao, cor) in enumerate(zip(nomes_cenarios, variacoes, cores)):
        y = 1.25 + i * 0.6
        
        # Linha do cenário
        adicionar_retangulo(slide, 0.3, y, 9.4, 0.55, COR_BRANCO, cor)
        
        # Dados
        dados = [
            cenario,
            variacao,
            formatar_moeda(cenarios[i]['receita_proj']),
            formatar_moeda(cenarios[i]['lucro_proj']),
            formatar_moeda(cenarios[i]['imposto_proj']),
            f"{(cenarios[i]['lucro_proj'] / cenarios[i]['receita_proj'] * 100):.1f}%"
        ]
        
        for j, (dado, larg) in enumerate(zip(dados, larguras)):
            x = 0.3 + sum(larguras[:j])
            alinhamento = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            cor_texto = cor if j == 0 else COR_PRETO
            
            adicionar_texto(slide, dado, x + 0.05, y + 0.1, larg, 0.35,
                            tamanho=9, cor=cor_texto, alinhamento=alinhamento, fonte="Arial")
    
    # Logo sobre o footer
    adicionar_logo_auditar(slide)
    
    # Texto do footer (posicionado acima da barra)
    adicionar_texto(slide, "🏢 AUDITAR - CONTABILIDADE CONSULTIVA 🏢", 0.3, 4.85, 9.4, 0.3,
                    tamanho=10, cor=COR_DOURADO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Footer dourado
    adicionar_retangulo(slide, 0, 5.2, 10, 0.1, COR_DOURADO)



# ============================================================
# SLIDE 4 - DRE COMPARATIVO ESTILO PDF
# ============================================================
def slide_dre_pdf(prs, dados_mensais, cenarios):
    """Slide DRE no estilo PDF com dourado e branco"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide, COR_BRANCO)
    
    # Header dourado
    adicionar_retangulo(slide, 0, 0, 10, 0.12, COR_DOURADO)
    adicionar_texto(slide, "04", 0.4, 0.02, 1, 0.6, tamanho=28, cor=COR_BRANCO, bold=True, fonte="Arial")
    adicionar_texto(slide, "📊 DRE COMPARATIVO ENTRE CENÁRIOS", 1.3, 0.02, 8, 0.6, tamanho=18, cor=COR_PRETO, bold=True, fonte="Arial")
    
    # Tabela comparativa estilo PDF
    cols = ["Cenário", "Receita", "Custos", "Lucro Op.", "Imposto", "Lucro Líq."]
    col_x = [0.3, 1.8, 3.3, 4.8, 6.3, 7.8]
    col_w = [1.4, 1.4, 1.4, 1.4, 1.4, 1.4]
    
    # Header da tabela
    adicionar_retangulo(slide, 0.3, 1.0, 9.4, 0.4, COR_CINZA_ESC)
    for i, (col, cx, cw) in enumerate(zip(cols, col_x, col_w)):
        adicionar_texto(slide, col, cx + 0.05, 1.05, cw, 0.3,
                        tamanho=8, cor=COR_BRANCO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
    
    # Dados
    receita_atual = sum(d["receita_bruta"] for d in dados_mensais)
    custos_atual = sum(d["custos"] for d in dados_mensais)
    lucro_atual = sum(d["lucro_operacional"] for d in dados_mensais)
    
    nomes_cenarios = ["ATUAL", "OTIMISTA", "REALISTA", "PESSIMISTA"]
    dados_cenarios = [
        (receita_atual, custos_atual, lucro_atual, lucro_atual * 0.06, lucro_atual * 0.94),
        (cenarios[0]['receita_proj'], cenarios[0]['custo_proj'], cenarios[0]['lucro_proj'], cenarios[0]['imposto_proj'], cenarios[0]['lucro_liq_proj']),
        (cenarios[1]['receita_proj'], cenarios[1]['custo_proj'], cenarios[1]['lucro_proj'], cenarios[1]['imposto_proj'], cenarios[1]['lucro_liq_proj']),
        (cenarios[2]['receita_proj'], cenarios[2]['custo_proj'], cenarios[2]['lucro_proj'], cenarios[2]['imposto_proj'], cenarios[2]['lucro_liq_proj']),
    ]
    
    for row_idx, (nome, dados) in enumerate(zip(nomes_cenarios, dados_cenarios)):
        y = 1.45 + row_idx * 0.4
        adicionar_retangulo(slide, 0.3, y, 9.4, 0.38, COR_BRANCO, COR_DOURADO)
        
        valores = [nome] + [formatar_moeda(v) for v in dados]
        for i, (val, cx, cw) in enumerate(zip(valores, col_x, col_w)):
            alinhamento = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            adicionar_texto(slide, val, cx + 0.05, y + 0.05, cw, 0.3,
                            tamanho=7, cor=COR_PRETO, alinhamento=alinhamento, fonte="Arial")
    
    # Footer dourado (primeiro, para ficar atrás)
    adicionar_retangulo(slide, 0, 5.2, 10, 0.1, COR_DOURADO)
    adicionar_texto(slide, "🏢 AUDITAR - CONTABILIDADE CONSULTIVA 🏢", 0.3, 5.18, 10, 0.4,
                    tamanho=10, cor=COR_BRANCO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Logo sobre o footer (depois, para ficar por cima)
    adicionar_logo_auditar(slide)

# ============================================================
# SLIDE 5 - RECOMENDAÇÕES ESTILIZADAS ESTILO PDF
# ============================================================
def slide_recomendacoes_pdf(prs, recomendacoes_ia, cenarios):
    """Slide de recomendações no estilo PDF com dourado e branco"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide, COR_BRANCO)

    # Header sem barra dourada
    adicionar_texto(slide, "05", 0.4, 0.15, 1, 0.6, tamanho=28, cor=COR_PRETO, bold=True, fonte="Arial")
    adicionar_texto(slide, "🎯 RECOMENDAÇÕES ESTRATÉGICAS", 1.3, 0.15, 8, 0.6, tamanho=18, cor=COR_PRETO, bold=True, fonte="Arial")

    # Área principal de recomendações sem borda dourada
    adicionar_retangulo(slide, 0.3, 1.2, 9.4, 3.5, COR_BRANCO, COR_CINZA_ESC)

    # Recomendações por cenário
    recomendacoes = [
        ("OTIMISTA", [
            "• Investir 20% do lucro em expansão",
            "• Contratar equipe especializada",
            "• Diversificar fontes de receita",
            "• Preparar estrutura para crescimento"
        ]),
        ("REALISTA", [
            "• Manter operações otimizadas",
            "• Focar em eficiência operacional",
            "• Buscar novas oportunidades",
            "• Reserva de contingência 10%"
        ]),
        ("PESSIMISTA", [
            "• Cortar custos não essenciais",
            "• Renegociar contratos e prazos",
            "• Focar em receitas recorrentes",
            "• Preparar plano de recuperação"
        ])
    ]

    for i, (titulo, recs) in enumerate(recomendacoes):
        x = 0.3 + (i % 2) * 4.7
        y = 1.1 + (i // 2) * 1.8

        # Card sem borda dourada
        adicionar_retangulo(slide, x, y, 4.6, 1.6, COR_BRANCO, COR_CINZA_ESC)

        # Header sem fundo dourado
        adicionar_texto(slide, titulo, x + 0.1, y + 0.05, 4.4, 0.2,
                        tamanho=9, cor=COR_PRETO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")

        # Recomendações
        for j, rec in enumerate(recs):
            adicionar_texto(slide, rec, x + 0.1, y + 0.4 + j * 0.25, 4.4, 0.2,
                            tamanho=7, cor=COR_PRETO, fonte="Arial")

    # Logo apenas no rodapé
    adicionar_logo_auditar(slide)

    # Footer sem barra dourada
    adicionar_texto(slide, "🏢 AUDITAR - CONTABILIDADE CONSULTIVA 🏢", 0.3, 5.2, 9.4, 0.3,
                    tamanho=10, cor=COR_PRETO, bold=True, alinhamento=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6 - PLANO DE AÇÃO ESTILO PDF
# ============================================================
def slide_plano_acao_pdf(prs, nome_empresa, cenarios):
    """Slide de plano de ação no estilo PDF com dourado e branco"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide, COR_BRANCO)
    
    # Header dourado
    adicionar_retangulo(slide, 0, 0, 10, 0.12, COR_DOURADO)
    adicionar_texto(slide, "06", 0.4, 0.02, 1, 0.6, tamanho=28, cor=COR_BRANCO, bold=True, fonte="Arial")
    adicionar_texto(slide, "⚡ PLANO DE AÇÃO ESTRATÉGICO", 1.3, 0.02, 8, 0.6, tamanho=18, cor=COR_PRETO, bold=True, fonte="Arial")
    
    # Timeline de ações estilo PDF
    meses = ["30 dias", "60 dias", "90 dias", "180 dias"]
    acoes = [
        ["Análise detalhada", "Definir KPIs", "Reestruturação", "Monitoramento"],
        ["Otimização fiscal", "Treinamento equipe", "Expansão", "Ajustes finos"],
        ["Implementação", "Processos", "Novos mercados", "Consolidação"],
        ["Avaliação", "Resultados", "Escala", "Sustentabilidade"]
    ]
    
    for i, (mes, acoes_mes) in enumerate(zip(meses, acoes)):
        x = 0.3 + i * 2.3
        y = 1.2
        
        # Card do período
        adicionar_retangulo(slide, x, y, 2.1, 3.0, COR_BRANCO, COR_DOURADO)
        
        # Header
        adicionar_retangulo(slide, x, y, 2.1, 0.3, COR_DOURADO)
        adicionar_texto(slide, mes, x + 0.1, y + 0.05, 1.9, 0.2,
                        tamanho=9, cor=COR_BRANCO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
        
        # Ações
        for j, acao in enumerate(acoes_mes):
            adicionar_texto(slide, f"• {acao}", x + 0.1, y + 0.4 + j * 0.6, 1.9, 0.5,
                            tamanho=7, cor=COR_PRETO, fonte="Arial")
    
    # KPIs de sucesso
    adicionar_retangulo(slide, 0.3, 4.4, 9.4, 0.4, COR_DOURADO)
    adicionar_texto(slide, "KPIs DE SUCESSO: Redução 15% custos | Aumento 25% eficiência | Economia R$ 1.5M tributos", 
                   0.5, 4.45, 8.4, 0.3, tamanho=9, cor=COR_BRANCO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
    
    # Logo sobre o footer
    adicionar_logo_auditar(slide)
    
    # Texto do footer (posicionado acima da barra)
    adicionar_texto(slide, "🏢 AUDITAR - CONTABILIDADE CONSULTIVA 🏢", 0.3, 4.85, 9.4, 0.3,
                    tamanho=10, cor=COR_DOURADO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Footer dourado
    adicionar_retangulo(slide, 0, 5.2, 10, 0.1, COR_DOURADO)

# ============================================================
# SLIDE 7 - COMPARAÇÃO ENTRE PERÍODOS ESTILO PDF
# ============================================================
def slide_comparacao_periodos_pdf(prs, dados_comparacao):
    """Slide de comparação entre períodos no estilo PDF"""
    if not dados_comparacao:
        return
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide, COR_BRANCO)
    
    # Header dourado
    adicionar_retangulo(slide, 0, 0, 10, 0.12, COR_DOURADO)
    adicionar_texto(slide, "07", 0.4, 0.02, 1, 0.6, tamanho=28, cor=COR_BRANCO, bold=True, fonte="Arial")
    adicionar_texto(slide, "📈 COMPARAÇÃO ENTRE PERÍODOS", 1.3, 0.02, 8, 0.6, tamanho=18, cor=COR_PRETO, bold=True, fonte="Arial")
    
    # Tabela de comparação
    cols = ["Mês", "Receita", "Var %", "Lucro Op.", "Var %", "Margem %", "Var %"]
    col_x = [0.3, 1.5, 2.8, 4.1, 5.4, 6.7, 8.0]
    col_w = [1.1, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2]
    
    # Header da tabela
    adicionar_retangulo(slide, 0.3, 0.8, 9.4, 0.4, COR_CINZA_ESC)
    for i, (col, cx, cw) in enumerate(zip(cols, col_x, col_w)):
        adicionar_texto(slide, col, cx + 0.05, 0.85, cw, 0.3,
                        tamanho=7, cor=COR_BRANCO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
    
    # Dados da comparação (espaçamento reduzido para caber 12 meses)
    for i, dado in enumerate(dados_comparacao['dados_mensais']):
        y = 1.25 + i * 0.16
        adicionar_retangulo(slide, 0.3, y, 9.4, 0.15, COR_BRANCO, COR_DOURADO)
        
        valores = [
            dado['mes'],
            formatar_moeda(dado['receita']),
            f"{dado['var_receita']:+.1f}%",
            formatar_moeda(dado['lucro']),
            f"{dado['var_lucro']:+.1f}%",
            f"{dado['margem']:.1f}%",
            f"{dado['var_margem']:+.1f}%"
        ]
        
        for j, (val, cx, cw) in enumerate(zip(valores, col_x, col_w)):
            alinhamento = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            cor_indicador = COR_PRETO
            if j in [2, 5, 6]:  # Colunas de variação
                if '+' in val:
                    cor_indicador = COR_DOURADO
                elif '-' in val:
                    cor_indicador = COR_CINZA_ESC
            
            adicionar_texto(slide, val, cx + 0.05, y + 0.02, cw, 0.15,
                            tamanho=5, cor=cor_indicador, alinhamento=alinhamento, fonte="Arial")
    
    # Tendência geral
    y_base = 1.25 + len(dados_comparacao['dados_mensais']) * 0.16 + 0.1
    adicionar_retangulo(slide, 0.3, y_base, 9.4, 0.35, COR_BRANCO, COR_DOURADO)
    
    texto_tendencia = f"Evolução ({dados_comparacao['primeiro_mes']} a {dados_comparacao['ultimo_mes']}): Receita {dados_comparacao['var_receita_total']:+.1f}% | Lucro {dados_comparacao['var_lucro_total']:+.1f}%"
    adicionar_texto(slide, texto_tendencia, 0.5, y_base + 0.1, 8.4, 0.3,
                    tamanho=8, cor=COR_PRETO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
    
    # Logo sobre o footer
    adicionar_logo_auditar(slide)
    
    # Texto do footer (posicionado acima da barra)
    adicionar_texto(slide, "🏢 AUDITAR - CONTABILIDADE CONSULTIVA 🏢", 0.3, 4.85, 9.4, 0.3,
                    tamanho=10, cor=COR_DOURADO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Footer dourado
    adicionar_retangulo(slide, 0, 5.2, 10, 0.1, COR_DOURADO)

# ============================================================
# SLIDE 8 - AGRADECIMENTOS ESTILO PDF
# ============================================================
def slide_agradecimentos_pdf(prs, nome_empresa):
    """Slide de agradecimentos no estilo PDF"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide, COR_BRANCO)
    
    # Header dourado
    adicionar_retangulo(slide, 0, 0, 10, 0.12, COR_DOURADO)
    adicionar_texto(slide, "08", 0.4, 0.02, 1, 0.6, tamanho=28, cor=COR_BRANCO, bold=True, fonte="Arial")
    
    # Mensagem principal de agradecimento
    adicionar_texto(slide, "MUITO OBRIGADO!", 0.3, 1.2, 9.4, 0.8,
                    tamanho=36, cor=COR_DOURADO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
    
    # Nome da empresa
    adicionar_texto(slide, f"{nome_empresa.upper()}", 0.3, 2.1, 9.4, 0.5,
                    tamanho=20, cor=COR_PRETO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
    
    # Linha separadora
    adicionar_retangulo(slide, 2.0, 2.7, 6.0, 0.02, COR_DOURADO)
    
    
    # Informações de contato   
    adicionar_texto(slide, "📱 (47) 3247-5001  ", 0.3, 4.9, 9.4, 0.3,
                    tamanho=11, cor=COR_DOURADO, bold=True, alinhamento=PP_ALIGN.CENTER, fonte="Arial")
    
    # Logo sobre o footer
    adicionar_logo_auditar(slide)
    
    # Texto do footer (posicionado acima da barra)
    adicionar_texto(slide, "🏢 AUDITAR - CONTABILIDADE CONSULTIVA 🏢", 0.3, 4.85, 9.4, 0.3,
                    tamanho=10, cor=COR_DOURADO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Footer dourado
    adicionar_retangulo(slide, 0, 5.2, 10, 0.1, COR_DOURADO)

# ============================================================
def gerar_cenarios_com_ia(dados_mensais, nome_empresa, bundle_dir):
    """Gera cenários otimista, realista e pessimista"""
    print("[IA] Gerando cenários econômicos...")
    
    receita_atual = sum(d["receita_bruta"] for d in dados_mensais)
    custos_atual = sum(d["custos"] for d in dados_mensais)
    
    cenarios = [
        {
            "nome": "Otimista",
            "receita_proj": receita_atual * 1.3,
            "custo_proj": custos_atual * 1.1,
            "lucro_proj": (receita_atual * 1.3) - (custos_atual * 1.1),
            "imposto_proj": (receita_atual * 1.3) * 0.06,
            "lucro_liq_proj": ((receita_atual * 1.3) - (custos_atual * 1.1)) - ((receita_atual * 1.3) * 0.06)
        },
        {
            "nome": "Realista", 
            "receita_proj": receita_atual * 1.1,
            "custo_proj": custos_atual * 1.05,
            "lucro_proj": (receita_atual * 1.1) - (custos_atual * 1.05),
            "imposto_proj": (receita_atual * 1.1) * 0.06,
            "lucro_liq_proj": ((receita_atual * 1.1) - (custos_atual * 1.05)) - ((receita_atual * 1.1) * 0.06)
        },
        {
            "nome": "Pessimista",
            "receita_proj": receita_atual * 0.8,
            "custo_proj": custos_atual * 0.95,
            "lucro_proj": (receita_atual * 0.8) - (custos_atual * 0.95),
            "imposto_proj": (receita_atual * 0.8) * 0.06,
            "lucro_liq_proj": ((receita_atual * 0.8) - (custos_atual * 0.95)) - ((receita_atual * 0.8) * 0.06)
        }
    ]
    
    print("[IA] Cenários gerados com sucesso!")
    return cenarios

def calcular_comparacao_periodos(dados_mensais):
    """Calcula dados de comparação entre períodos"""
    if len(dados_mensais) < 2:
        return None
    
    comparacao = []
    receita_anterior = None
    lucro_anterior = None
    margem_anterior = None
    
    for i, d in enumerate(dados_mensais):
        receita_atual = d['receita_bruta']
        lucro_atual = d['lucro_operacional']
        margem_atual = (lucro_atual / receita_atual * 100) if receita_atual > 0 else 0
        
        # Calcular variações
        if i > 0 and receita_anterior:
            var_receita = ((receita_atual - receita_anterior) / receita_anterior * 100) if receita_anterior > 0 else 0
        else:
            var_receita = 0
        
        if i > 0 and lucro_anterior:
            var_lucro = ((lucro_atual - lucro_anterior) / lucro_anterior * 100) if lucro_anterior > 0 else 0
        else:
            var_lucro = 0
        
        if i > 0 and margem_anterior:
            var_margem = margem_atual - margem_anterior
        else:
            var_margem = 0
        
        # Mapear número do mês para nome
        meses_nomes = ['Jan', 'Fev', 'Mar', 'Abr', 'Mai', 'Jun', 'Jul', 'Ago', 'Set', 'Out', 'Nov', 'Dez']
        mes_num = int(d['mes']) if isinstance(d['mes'], str) else d['mes']
        mes_nome = meses_nomes[mes_num - 1] if 1 <= mes_num <= 12 else str(mes_num)
        
        comparacao.append({
            'mes': mes_nome,
            'receita': receita_atual,
            'var_receita': var_receita,
            'lucro': lucro_atual,
            'var_lucro': var_lucro,
            'margem': margem_atual,
            'var_margem': var_margem
        })
        
        receita_anterior = receita_atual
        lucro_anterior = lucro_atual
        margem_anterior = margem_atual
    
    # Calcular tendência geral
    primeiro = dados_mensais[0]
    ultimo = dados_mensais[-1]
    
    var_receita_total = ((ultimo['receita_bruta'] - primeiro['receita_bruta']) / primeiro['receita_bruta'] * 100) if primeiro['receita_bruta'] > 0 else 0
    var_lucro_total = ((ultimo['lucro_operacional'] - primeiro['lucro_operacional']) / primeiro['lucro_operacional'] * 100) if primeiro['lucro_operacional'] > 0 else 0
    
    return {
        'dados_mensais': comparacao,
        'var_receita_total': var_receita_total,
        'var_lucro_total': var_lucro_total,
        'primeiro_mes': primeiro['mes'],
        'ultimo_mes': ultimo['mes']
    }

def analisar_com_ia_pdf(dados_mensais, nome_empresa, bundle_dir):
    """Análise avançada com IA para estilo PDF"""
    print("[IA] Analisando dados com Inteligência Artificial...")

    receita_t = sum(d["receita_bruta"] for d in dados_mensais)
    custos_t = sum(d["custos"] for d in dados_mensais)
    desp_t = sum(d["despesas"] for d in dados_mensais)
    lucro_t = sum(d["lucro_operacional"] for d in dados_mensais)

    prompt = f"""Você é um especialista em contabilidade e planejamento tributário brasileiro.

EMPRESA: {nome_empresa}
RECEITA TOTAL: R$ {receita_t:,.2f}
CUSTOS: R$ {custos_t:,.2f}
DESPESAS: R$ {desp_t:,.2f}
LUCRO: R$ {lucro_t:,.2f}

Forneça uma análise estratégica completa:
1. Resumo da performance financeira
2. Principais oportunidades de otimização
3. Riscos e recomendações
4. Sugestões de planejamento tributário

Responda de forma profissional e concisa:"""

    try:
        import anthropic
        api_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "data", "api_config.json")
        
        with open(api_path, "r", encoding="utf-8") as f:
            api_config = json.load(f)

        api_key = api_config.get("api_key", "")
        model = api_config.get("model", "claude-sonnet-4-6")

        if not api_key:
            raise ValueError("API KEY não encontrada")

        client = anthropic.Anthropic(api_key=api_key)
        message = client.messages.create(model=model, max_tokens=1024, messages=[{"role": "user", "content": prompt}])
        resposta = message.content[0].text.strip()

        print("[IA] Análise concluída com sucesso!")
        return resposta

    except Exception as e:
        print(f"[IA] Aviso: usando análise padrão ({e})")
        return f"A empresa {nome_empresa} apresenta receita de {formatar_moeda(receita_t)} com margem de {(lucro_t/receita_t*100):.1f}%. Recomenda-se focar em otimização tributária e controle de custos para maximizar a lucratividade."

# ============================================================
# FUNÇÃO PRINCIPAL V2.0 ESTILO PDF
# ============================================================
def gerar_apresentacao_pptx_pdf(
    dados_mensais,
    nome_empresa,
    responsavel,
    bundle_dir,
    api_key,
    cores_personalizadas=None
):
    """
    Gera apresentação PowerPoint V2.0 no estilo PDF profissional
    """
    if not dados_mensais:
        print("[ERRO] Nenhum dado disponível para gerar a apresentação.")
        return

    # Usar cores personalizadas se fornecidas
    if cores_personalizadas:
        atualizar_cores(cores_personalizadas)
        print("🎨 Cores personalizadas aplicadas!")

    print("\n" + "=" * 60)
    print("GERADOR DE APRESENTAÇÃO V2.0 - ESTILO PDF PROFISSIONAL")
    print("=" * 60)

    # Análises com IA
    analise_ia = analisar_com_ia_pdf(dados_mensais, nome_empresa, bundle_dir)
    cenarios = gerar_cenarios_com_ia(dados_mensais, nome_empresa, bundle_dir)
    dados_comparacao = calcular_comparacao_periodos(dados_mensais)

    # Criar apresentação
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    data_emissao = datetime.now().strftime("%d/%m/%Y")

    print("[PPTX PDF] Gerando slides estilo profissional com estrutura V2.0...")
    slide_capa_pdf(prs, nome_empresa, responsavel, data_emissao)
    slide_sumario_pdf(prs, dados_mensais, nome_empresa)
    slide_cenarios_pdf(prs, dados_mensais, cenarios)
    slide_dre_pdf(prs, dados_mensais, cenarios)
    slide_recomendacoes_pdf(prs, analise_ia, cenarios)
    slide_plano_acao_pdf(prs, nome_empresa, cenarios)
    
    # Adicionar slide de comparação se houver dados suficientes
    if dados_comparacao:
        slide_comparacao_periodos_pdf(prs, dados_comparacao)
    
    slide_agradecimentos_pdf(prs, nome_empresa)

    # Salvar arquivo
    nome_arquivo = f"Apresentacao_EstiloPDF_{nome_empresa.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    caminho = os.path.join(bundle_dir, nome_arquivo)
    prs.save(caminho)

    logo_existe = os.path.exists(os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo_auditar.png"))
    
    print("\n[OK] Apresentação estilo PDF gerada com sucesso!")
    print(f"     Arquivo: {nome_arquivo}")
    print(f"     Local: {caminho}")
    print(f"     Estilo: Profissional PDF | Logo: {'SIM' if logo_existe else 'NÃO'}")
    return caminho




