"""
GERADOR INTELIGENTE DE APRESENTAÇÕES COM IA LOCAL
Integração com Ollama para criação dinâmica de slides via linguagem natural
"""

import os
import sys
import json
import re
import httpx
import subprocess
import tempfile
from typing import List, Dict, Any, Optional
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Configurações
OLLAMA_HOST = "http://localhost:11434"
DEFAULT_MODEL = "tinyllama"


class GeradorIAInteligente:
    """Gerador de apresentações usando IA local (Ollama)"""
    
    def __init__(self, model: str = DEFAULT_MODEL):
        self.model = model
        self.ollama_available = self._check_ollama()
        self.client = httpx.AsyncClient(timeout=300.0)
        
    def _check_ollama(self) -> bool:
        """Verifica se Ollama está rodando"""
        try:
            import urllib.request
            urllib.request.urlopen(f"{OLLAMA_HOST}/api/tags", timeout=2)
            print("Ollama detectado")
            return True
        except:
            print("Ollama nao disponivel - usando fallback")
            return False
    
    async def gerar_apresentacao_inteligente(
        self,
        dados_mensais: List[Dict],
        nome_empresa: str,
        responsavel: str,
        comando_estilo: str,
        num_slides: int = 6
    ) -> str:
        """
        Gera apresentação baseada em comando em linguagem natural
        
        Exemplos de comandos:
        - "Crie uma apresentação com fundo azul marinho, títulos em dourado, estilo corporativo elegante"
        - "Faça slides modernos com fundo escuro, acentos em verde neon, layout minimalista"
        - "Apresentação clean com branco e cinza, tipografia profissional, ícones sutis"
        """
        
        print(f"\nAnalisando comando: '{comando_estilo}'")
        print(f"Dados: {len(dados_mensais)} meses de {nome_empresa}")
        
        # 1. Extrair requisitos de design via IA
        requisitos_design = await self._analisar_requisitos(comando_estilo)
        print(f"Design escolhido: {requisitos_design.get('nome_estilo', 'Personalizado')}")
        
        # 2. Gerar código Python dinâmico para os slides
        codigo_gerado = await self._gerar_codigo_slides(
            dados_mensais, 
            nome_empresa, 
            responsavel,
            requisitos_design,
            num_slides
        )
        
        # 3. Executar código gerado
        print("Construindo apresentacao...")
        filepath = await self._executar_codigo_gerado(
            codigo_gerado,
            dados_mensais,
            nome_empresa,
            responsavel,
            requisitos_design
        )
        
        return filepath

    async def _call_ollama_chat(self, messages: List[Dict[str, str]]) -> str:
        """Faz chamada à API de Chat do Ollama"""
        try:
            payload = {
                "model": self.model,
                "messages": messages,
                "stream": False,
                "options": {"temperature": 0.4, "num_predict": 2048, "num_ctx": 4096}
            }
            
            response = await self.client.post(
                f"{OLLAMA_HOST}/api/chat",
                json=payload,
                timeout=300.0
            )
            response.raise_for_status()
            data = response.json()
            return data.get("message", {}).get("content", "")
        except Exception as e:
            print(f"[ERRO] Chat API: {type(e).__name__} - {e}")
            raise

    async def analisar_financas(
        self,
        dados_mensais: List[Dict],
        nome_empresa: str,
        comando_personalizado: str = ""
    ) -> str:
        """Analisa os dados financeiros e retorna um diagnóstico em texto"""
        
        if not self.ollama_available:
            return "Ollama não disponível. Verifique se o serviço está rodando."

        # Preparar dados para o prompt
        resumo = []
        for d in dados_mensais[-6:]: 
            resumo.append(f"Mês {d['mes']}/{d['ano']}: Receita R$ {d['receita_bruta']:,.2f}, Lucro R$ {d['lucro_operacional']:,.2f}")
        
        dados_texto = "\n".join(resumo)
        
        instrucao = f"Você é um consultor financeiro brasileiro. Analise estes dados da empresa {nome_empresa} e forneça um diagnóstico estratégico curto (máximo 2 parágrafos) com 1 tendência e 2 sugestões estratégicas."
        instrucao += "\nNÃO explique o significado dos termos (ex: não explique o que é redução de custos), apenas dê as sugestões práticas."
        if comando_personalizado:
            instrucao += f"\nComando extra: {comando_personalizado}"

        # Template rígido para TinyLlama
        prompt_final = f"### Instruction:\n{instrucao}\n\nDADOS:\n{dados_texto}\n\n### Response:\n"

        try:
            # Tentar Chat API primeiro (corrigida)
            try:
                print(f"Analisando {nome_empresa} via Chat API...")
                messages = [{"role": "user", "content": prompt_final}]
                analise = await self._call_ollama_chat(messages)
            except:
                print("Chat API falhou, usando Generate...")
                analise = await self._call_ollama(prompt_final)
            
            if not analise:
                return "Erro: IA retornou vazio."

            resposta_limpa = analise.strip()
            
            # Filtro de prefixos e limpeza de identidade (AGRESSIVO)
            padrões_remover = [
                "Você é um", "Como consultor", "Sou um consultor",
                "Ao analisar", "Diagnóstico:", "Resposta:", "### Response:",
                "Com o CNPJ", "financeiro brasileiro", "analizei", "analisamos",
                "esteos dações", "dações do grupo", "proporciona um diagnóstico",
                "Comando extra:", "Comando adicional:", "DAO:", "DADOS:", "Significa identificar"
            ]
            
            linhas = resposta_limpa.split('\n')
            final_linhas = []
            for i, linha in enumerate(linhas):
                l_lower = linha.lower()
                # Pula linhas de instrução ou definições inúteis nas primeiras linhas
                if i < 5 and any(p.lower() in l_lower for p in padrões_remover):
                    continue
                final_linhas.append(linha)
            
            resposta_limpa = "\n".join(final_linhas).strip()

            # Se ainda começar com algo estranho, buscar o primeiro ponto de dado real
            for foco in ["A empresa", "Com base", "Observa-se", "Analisando", "1. ", "Tendência", "Sugestão"]:
                pos = resposta_limpa.find(foco)
                if 0 <= pos < 250: 
                    resposta_limpa = resposta_limpa[pos:].strip()
                    break

            return resposta_limpa
            
        except Exception as e:
            print(f"[ERRO] Falha na análise: {e}")
            return f"Erro na análise: {str(e)}"
    
    async def _analisar_requisitos(self, comando: str) -> Dict[str, Any]:
        """Usa IA para extrair requisitos de design do comando do usuário"""

        print(f"Iniciando analise de requisitos para: '{comando}'")

        if not self.ollama_available:
            print("Ollama nao disponivel, usando fallback")
            return self._design_fallback(comando)

        try:
            print("Chamando Ollama (Chat API)...")
            messages = [
                {"role": "system", "content": "Voce e um designer. Responda APENAS o JSON com cores RGB baseadas no comando."},
                {"role": "user", "content": f'Command: "{comando}". Fill this JSON with colors: {{"nome_estilo":"","cores":{{"fundo":[255,255,255],"titulo":[0,0,0],"texto":[0,0,0],"acento":[0,0,0],"destaque":[0,0,0]}},"fontes":{{"titulo":"Calibri","corpo":"Calibri"}},"estilo_visual":"","elementos_especiais":[],"layout":"","descricao":""}}'}
            ]
            response = await self._call_ollama_chat(messages)
            print(f"Resposta recebida (tamanho: {len(response)} chars)")

            # Substituição manual de cores óbvias se a IA falhar (Camada Extra de Segurança)
            cmd_l = comando.lower()
            if "vermelho" in cmd_l or "red" in cmd_l:
                response = response.replace("[0,0,0]", "[220,53,69]").replace("[0, 0, 0]", "[220,53,69]")
            elif "azul" in cmd_l or "blue" in cmd_l:
                response = response.replace("[0,0,0]", "[0,112,192]").replace("[0, 0, 0]", "[0,112,192]")
            elif "verde" in cmd_l or "green" in cmd_l:
                response = response.replace("[0,0,0]", "[40,167,69]").replace("[0, 0, 0]", "[40,167,69]")

            # Limpar possiveis letras [R, G, B] que o modelo possa ter repetido
            response = re.sub(r'[RGB]', '0', response)

            # Extrair JSON da resposta
            json_str = self._extrair_json(response)
            print(f"JSON extraido (tamanho: {len(json_str)} chars)")

            # Tentar fazer parsing do JSON
            requisitos = None
            try:
                # Definir a string de limpeza baseada no que foi extraido
                json_str_clean = json_str
                
                # Limpeza bruta: transforma [[R,G,B], ...] em apenas [R,G,B]
                # Pega o primeiro conjunto de 3 números dentro de colchetes aninhados
                json_str_clean = re.sub(r'\[\s*\[\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\].*?\]', r'[\1, \2, \3]', json_str_clean)
                # Remove colchetes triplos ou quádruplos se existirem
                json_str_clean = json_str_clean.replace("[[[", "[").replace("]]]", "]")
                json_str_clean = json_str_clean.replace("[[", "[").replace("]]", "]")
                
                try:
                    requisitos = json.loads(json_str_clean)
                except:
                    # Tentativa extra: ast.literal_eval (mais leniente com aspas simples)
                    try:
                        import ast
                        requisitos = ast.literal_eval(json_str_clean)
                        print("JSON recuperado via ast.literal_eval")
                    except:
                        # Tentar encontrar blocos de cores via regex se tudo falhar
                        cores_encontradas = re.findall(r'\[(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\]', json_str)
                    if len(cores_encontradas) >= 3:
                        requisitos = self._design_fallback(comando)
                        keys = ["fundo", "titulo", "texto", "acento", "destaque"]
                        for i, cor in enumerate(cores_encontradas[:5]):
                            requisitos["cores"][keys[i]] = [int(c) for c in cor]
                
                if requisitos:
                    print("JSON extraido/recuperado com sucesso")
            except Exception as e:
                print(f"Erro no parsing inicial: {e}")

            if requisitos is None:
                print("[AVISO] Nao foi possivel extrair JSON, usando fallback inteligente")
                return self._design_fallback(comando)

            # Validar e corrigir chaves mal escritas (ex: "core" vs "cores")
            if "core" in requisitos and "cores" not in requisitos:
                requisitos["cores"] = requisitos.pop("core")
            
            # Validar estrutura final
            try:
                requisitos = self._validar_requisitos(requisitos, comando)
                print("[OK] Requisitos validados")
            except Exception as e:
                print(f"[AVISO] Erro na validacao: {e}")
                return self._design_fallback(comando)

            return requisitos

        except Exception as e:
            print(f"[ERRO] Erro na anlise IA: {e}")
            import traceback
            traceback.print_exc()
            return self._design_fallback(comando)
    
    async def _gerar_codigo_slides(
        self,
        dados_mensais: List[Dict],
        nome_empresa: str,
        responsavel: str,
        design: Dict,
        num_slides: int
    ) -> str:
        """Gera código Python dinâmico para criar os slides"""
        
        # Formatar mês/ano com tratamento de tipo
        def fmt_mes_ano(mes_val, ano_val, mes_default='01', ano_default='2024'):
            try:
                mes = int(mes_val)
                mes_str = f"{mes:02d}"
            except:
                mes_str = str(mes_val).zfill(2)[:2] if mes_val else mes_default
            try:
                ano = int(ano_val)
                ano_str = str(ano)
            except:
                ano_str = str(ano_val) if ano_val else ano_default
            return f"{mes_str}/{ano_str}"
        
        # Calcular indicadores
        receita_total = sum(d["receita_bruta"] for d in dados_mensais)
        lucro_total = sum(d["lucro_operacional"] for d in dados_mensais)
        margem = (lucro_total / receita_total * 100) if receita_total > 0 else 0
        
        # Criar resumo dos dados
        resumo_dados = {
            "meses": len(dados_mensais),
            "receita_total": f"R$ {receita_total:,.2f}",
            "lucro_total": f"R$ {lucro_total:,.2f}",
            "margem": f"{margem:.1f}%",
            "primeiro_mes": fmt_mes_ano(dados_mensais[0].get('mes'), dados_mensais[0].get('ano')) if dados_mensais else "",
            "ultimo_mes": fmt_mes_ano(dados_mensais[-1].get('mes'), dados_mensais[-1].get('ano')) if dados_mensais else ""
        }
        
        prompt_codigo = f"""Você é um programador Python especialista em python-pptx.
Gere código Python para criar uma apresentação PowerPoint profissional.

REQUISITOS DE DESIGN:
- Estilo: {design.get('nome_estilo', 'Personalizado')}
- Fundo: RGB{design['cores']['fundo']} - APLIQUE EXATAMENTE ESTA COR AO FUNDO DE TODOS OS SLIDES
- Títulos: RGB{design['cores']['titulo']} - APLIQUE EXATAMENTE ESTA COR AOS TÍTULOS
- Textos: RGB{design['cores']['texto']} - APLIQUE EXATAMENTE ESTA COR AOS TEXTOS CORPO
- Acentos: RGB{design['cores']['acento']} - APLIQUE EXATAMENTE ESTA COR AOS ELEMENTOS DE DESTAQUE
- Elementos: {', '.join(design.get('elementos_especiais', []))}
- Layout: {design.get('layout', 'profissional')}

INSTRUÇÕES CRÍTICAS DE CONTRASTE:
- Se fundo for branco [255, 255, 255], use texto preto [0, 0, 0] ou cinza escuro
- Se fundo for escuro, use texto branco [255, 255, 255] ou claro
- NUNCA use texto escuro sobre fundo escuro
- NUNCA use texto claro sobre fundo claro
- Garanta legibilidade máxima

INSTRUÇÕES DE LAYOUT DE CARDS:
- NÃO use barras coloridas dentro dos cards de indicadores (slide 2 - RESUMO EXECUTIVO)
- NÃO use barras coloridas dentro dos cards de cenários (slide 4 - CENÁRIOS)
- NÃO use barras coloridas dentro dos cards de recomendações (slide 5 - RECOMENDAÇÕES)
- Cards devem ter fundo branco ou cinza claro, sem bordas coloridas
- Use apenas bordas cinza claras para os cards
- Mantenha barras coloridas no cabeçalho da tabela de DADOS MENSAIS (slide 3)

DADOS DA APRESENTAÇÃO:
- Empresa: {nome_empresa}
- Responsável: {responsavel}
- Período: {resumo_dados['primeiro_mes']} a {resumo_dados['ultimo_mes']}
- Receita Total: {resumo_dados['receita_total']}
- Lucro Total: {resumo_dados['lucro_total']}
- Margem: {resumo_dados['margem']}

ESTRUTURA DOS {num_slides} SLIDES:
1. CAPA - Título impactante, nome empresa, responsável, data
2. RESUMO EXECUTIVO - Dashboard com indicadores principais em CARDS (sem barras coloridas nos cards)
3. ANÁLISE FINANCEIRA - Gráfico visual da receita e lucro
4. CENÁRIOS - Projeções otimista, realista, pessimista em CARDS (sem barras coloridas nos cards)
5. RECOMENDAÇÕES - Estratégias baseadas nos dados em CARDS (sem barras coloridas nos cards)
6. CONCLUSÃO - Próximos passos e agradecimento

Gere código Python COMPLETO e EXECUTÁVEL que:
1. Importe todas as bibliotecas necessárias
2. Defina as cores RGB exatas
3. Crie funções auxiliares (adicionar_fundo, adicionar_texto, etc.)
4. Crie uma função principal gerar_apresentacao() que retorna o caminho do arquivo
5. Use python-pptx (Presentation, Inches, Pt, RGBColor, PP_ALIGN, MSO_SHAPE)
6. Aplique o design especificado fielmente
7. Salve com nome único baseado na data

IMPORTANTE - REGRAS DE CÓDIGO:
- Use APENAS nomes de variáveis em INGLÊS (ex: background_color, title_text, not fundo, titulo)
- NÃO use palavras em português como nomes de variáveis ou funções
- NÃO use aspas simples dentro de strings (use \" ou escape)
- Inclua tratamento de erro try/except em cada função auxiliar
- Teste se objetos existem antes de acessar propriedades
- Use nomes descritivos: slide, shape, text_frame, paragraph
- CADA função deve ser INDEPENDENTE - não use variáveis de outras funções
- Defina TODAS as variáveis necessárias DENTRO de cada função ou passe como parâmetro

EXEMPLO DE CÓDIGO CORRETO (copie este padrão):
```python
def slide_capa(prs, company_name, responsible):
    # Todas as variáveis DEFINIDAS AQUI dentro
    background_color = RGBColor(30, 58, 138)
    title_text = company_name
    # ... código do slide ...

def slide_dados(prs, summary_data):
    # Todas as variáveis DEFINIDAS AQUI dentro
    revenue = summary_data["receita_total"]
    # ... código do slide ...

# Função principal com TODOS os dados passados como parâmetros
def gerar_apresentacao():
    prs = Presentation()
    company_name = "{nome_empresa}"
    responsible = "{responsavel}"
    summary_data = {dict(resumo_dados)}
    slide_capa(prs, company_name, responsible)
    slide_dados(prs, summary_data)
    prs.save("arquivo.pptx")
```

O código deve estar pronto para execução imediata. Use apenas bibliotecas padrão + python-pptx.

CÓDIGO PYTHON:"""

        if not self.ollama_available:
            return self._codigo_fallback(design, nome_empresa, responsavel, resumo_dados, dados_mensais)
        
        try:
            messages = [
                {"role": "system", "content": "Voce e um programador Python expert em python-pptx. Gere APENAS codigo Python executavel, sem explicacoes."},
                {"role": "user", "content": prompt_codigo}
            ]
            codigo = await self._call_ollama_chat(messages)
            
            # Limpar código (remover markdown se presente)
            codigo = self._limpar_codigo(codigo)
            
            # Validar código básico
            if "def gerar_apresentacao" not in codigo:
                raise ValueError("Código não contém função principal")
            
            return codigo
            
        except Exception as e:
            print(f"[AVISO]  Erro na gerao de cdigo: {e}")
            return self._codigo_fallback(design, nome_empresa, responsavel, resumo_dados, dados_mensais)
    
    async def _executar_codigo_gerado(
        self,
        codigo: str,
        dados_mensais: List[Dict],
        nome_empresa: str,
        responsavel: str,
        design: Dict
    ) -> str:
        """Executa o código Python gerado de forma segura"""
        
        # Criar arquivo temporário para o código
        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False, encoding='utf-8') as f:
            f.write(codigo)
            temp_file = f.name
        
        try:
            # Preparar namespace seguro
            namespace = {
                'os': os,
                'sys': sys,
                'json': json,
                'datetime': datetime,
                '__name__': '__main__'
            }
            
            # Adicionar imports necessários ao namespace
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN
                from pptx.enum.shapes import MSO_SHAPE
                
                namespace.update({
                    'Presentation': Presentation,
                    'Inches': Inches,
                    'Pt': Pt,
                    'RGBColor': RGBColor,
                    'PP_ALIGN': PP_ALIGN,
                    'MSO_SHAPE': MSO_SHAPE
                })
            except ImportError:
                raise ImportError("python-pptx não instalado")
            
            # Executar código
            exec(codigo, namespace)
            
            # Chamar função principal
            if 'gerar_apresentacao' in namespace:
                resultado = namespace['gerar_apresentacao']()
                
                # O resultado pode ser o caminho ou None
                if resultado and os.path.exists(resultado):
                    return resultado
                
                # Procurar arquivo gerado
                return self._encontrar_arquivo_gerado(nome_empresa)
            else:
                raise ValueError("Função gerar_apresentacao não encontrada no código")
                
        except Exception as e:
            print(f"Erro ao executar codigo gerado: {e}")
            import traceback
            traceback.print_exc()
            print("Usando gerador fallback local...")
            
            # Verificar se algum arquivo foi criado mesmo com erro
            arquivo_recente = self._encontrar_arquivo_recente(nome_empresa, tempo_limite_segundos=10)
            if arquivo_recente:
                print(f"Arquivo encontrado apesar do erro: {arquivo_recente}")
                return arquivo_recente
            
            # Se não encontrou, gerar com fallback
            # Se design for string, usar _design_fallback para criar dict
            if isinstance(design, str):
                design = self._design_fallback(design)
            return self._gerar_fallback_local(
                dados_mensais, nome_empresa, responsavel, design
            )
        
        finally:
            # Limpar arquivo temporário
            try:
                os.unlink(temp_file)
            except:
                pass
    
    async def _call_ollama(self, prompt: str) -> str:
        """Faz chamada à API do Ollama"""
        
        try:
            payload = {
                "model": self.model,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "temperature": 0.7,
                    "num_predict": 2048,
                    "num_ctx": 4096
                }
            }
            
            response = await self.client.post(
                f"{OLLAMA_HOST}/api/generate",
                json=payload,
                headers={"Content-Type": "application/json"},
                timeout=300.0
            )
            
            response.raise_for_status()
            data = response.json()
            return data.get("response", "")
        
        except Exception as e:
            print(f"[ERRO] Erro ao chamar Ollama: {e}")
            raise
    
    def _extrair_json(self, texto: str) -> str:
        """Extrai JSON de texto que pode conter markdown ou outros caracteres"""

        # Tentar extrair bloco JSON de markdown
        if "```json" in texto:
            inicio = texto.find("```json") + 7
            fim = texto.find("```", inicio)
            if fim == -1:
                fim = len(texto)
            json_str = texto[inicio:fim].strip()
        elif "```" in texto:
            inicio = texto.find("```") + 3
            fim = texto.find("```", inicio)
            if fim == -1:
                fim = len(texto)
            json_str = texto[inicio:fim].strip()
        else:
            # Tentar encontrar primeiro { e último }
            inicio = texto.find("{")
            fim = texto.rfind("}")
            if inicio != -1 and fim != -1:
                json_str = texto[inicio:fim+1]
            else:
                json_str = texto

        # CORREÇÃO 1: Remover comentários (se houver)
        json_str = re.sub(r'//.*', '', json_str)
        json_str = re.sub(r'/\*.*?\*/', '', json_str, flags=re.DOTALL)

        # CORREÇÃO 2: Corrigir aspas simples em KEYS (padrão: 'key': -> "key":)
        json_str = re.sub(r"\'([^\']+)\'\s*:", r'"\1":', json_str)

        # CORREÇÃO 3: Corrigir chaves sem aspas (padrão: key: -> "key":)
        json_str = re.sub(r'(\s+)(\w+):', r'\1"\2":', json_str)

        # CORREÇÃO 4: Corrigir aspas simples em valores string (padrão: : 'value' -> : "value")
        json_str = re.sub(r':\s*\'([^\']+)\'', r': "\1"', json_str)

        # CORREÇÃO 5: Remover trailing commas
        json_str = re.sub(r',\s*}', '}', json_str)
        json_str = re.sub(r',\s*]', ']', json_str)

        # CORREÇÃO 6: Corrigir valores Python para JSON
        json_str = re.sub(r'True', 'true', json_str)
        json_str = re.sub(r'False', 'false', json_str)
        json_str = re.sub(r'None', 'null', json_str)

        # CORREÇÃO 7: Remover vírgulas duplicadas
        json_str = re.sub(r',\s*,', ',', json_str)

        # CORREÇÃO 8: Corrigir espaços entre propriedade e valor
        json_str = re.sub(r'(\w+)\s+([{\[])', r'\1\2', json_str)

        # CORREÇÃO 9: Remover caracteres de controle problemáticos
        json_str = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f]', '', json_str)

        # CORREÇÃO 10: Corrigir RGB sem colchetes (ex: 0, 0, 0 -> [0, 0, 0])
        json_str = re.sub(r'"fundo":\s*(\d+),\s*(\d+),\s*(\d+)', r'"fundo": [\1, \2, \3]', json_str)
        json_str = re.sub(r'"titulo":\s*(\d+),\s*(\d+),\s*(\d+)', r'"titulo": [\1, \2, \3]', json_str)
        json_str = re.sub(r'"texto":\s*(\d+),\s*(\d+),\s*(\d+)', r'"texto": [\1, \2, \3]', json_str)
        json_str = re.sub(r'"acento":\s*(\d+),\s*(\d+),\s*(\d+)', r'"acento": [\1, \2, \3]', json_str)
        json_str = re.sub(r'"destaque":\s*(\d+),\s*(\d+),\s*(\d+)', r'"destaque": [\1, \2, \3]', json_str)

        return json_str.strip()
    
    def _limpar_codigo(self, codigo: str) -> str:
        """Limpa código removendo markdown e formatação"""
        
        # Remover markdown
        if "```python" in codigo:
            codigo = codigo.split("```python")[1]
        elif "```" in codigo:
            codigo = codigo.split("```")[1] if codigo.startswith("```") else codigo
        
        # Remover explicações antes do código
        linhas = codigo.split('\n')
        inicio_codigo = 0
        
        for i, linha in enumerate(linhas):
            if linha.strip().startswith(('import ', 'from ', 'def ', '#')):
                inicio_codigo = i
                break
        
        return '\n'.join(linhas[inicio_codigo:]).strip()
    
    def _validar_requisitos(self, req: Dict, comando: str = "") -> Dict:
        """Valida e completa requisitos de design"""

        # Detectar se usuario quer preto explicitamente
        comando_lower = comando.lower()
        usar_preto = "preto" in comando_lower or "black" in comando_lower or "texto preto" in comando_lower
        usar_branco = "branco" in comando_lower or "white" in comando_lower or "fundo branco" in comando_lower

        padrao = {
            "nome_estilo": "Corporativo",
            "cores": {
                "fundo": [255, 255, 255],
                "titulo": [0, 0, 0] if usar_preto else [45, 45, 45],
                "texto": [0, 0, 0] if usar_preto else [80, 80, 80],
                "acento": [184, 143, 0],
                "destaque": [0, 112, 192]
            },
            "fontes": {
                "titulo": "Calibri",
                "corpo": "Calibri"
            },
            "estilo_visual": "profissional",
            "elementos_especiais": ["linhas decorativas"],
            "layout": "corporativo",
            "descricao": "Design corporativo clássico"
        }

        # Forçar fundo branco se solicitado
        if usar_branco:
            padrao["cores"]["fundo"] = [255, 255, 255]
            padrao["cores"]["titulo"] = [0, 0, 0]
            padrao["cores"]["texto"] = [0, 0, 0]
        
        # Mesclar com padrão (nível raiz)
        for key, value in padrao.items():
            if key not in req or not req[key]:
                req[key] = value
        
        # Garantir que 'cores' exista e tenha todas as chaves necessárias
        if "cores" not in req or not isinstance(req["cores"], dict):
            req["cores"] = padrao["cores"].copy()
        else:
            # Mesclar sub-chaves de cores
            for cor_key, cor_val in padrao["cores"].items():
                if cor_key not in req["cores"] or not req["cores"][cor_key]:
                    req["cores"][cor_key] = cor_val
        
        # Validar formato RGB de cada cor
        for cor_key, padrao_cor in padrao["cores"].items():
            cor = req["cores"].get(cor_key)
            if not isinstance(cor, list) or len(cor) != 3:
                req["cores"][cor_key] = padrao_cor
        
        return req
    
    def _design_fallback(self, comando: str) -> Dict:
        """Retorna design baseado em palavras-chave do comando"""

        comando_lower = comando.lower()
        usar_preto = "preto" in comando_lower or "black" in comando_lower or "texto preto" in comando_lower

        # Detectar cores
        if "azul" in comando_lower or "marinho" in comando_lower:
            if "dourado" in comando_lower or "gold" in comando_lower:
                return {
                    "nome_estilo": "Azul Marinho & Dourado",
                    "cores": {
                        "fundo": [25, 40, 65],
                        "titulo": [212, 175, 55],
                        "texto": [220, 220, 220],
                        "acento": [212, 175, 55],
                        "destaque": [100, 149, 237]
                    },
                    "fontes": {"titulo": "Calibri", "corpo": "Calibri"},
                    "estilo_visual": "luxuoso",
                    "elementos_especiais": ["gradientes sutis", "bordas douradas"],
                    "layout": "elegante",
                    "descricao": "Design luxuoso azul marinho com acentos dourados"
                }
            else:
                return {
                    "nome_estilo": "Azul Corporativo",
                    "cores": {
                        "fundo": [240, 248, 255],
                        "titulo": [0, 0, 0] if usar_preto else [0, 51, 102],
                        "texto": [0, 0, 0] if usar_preto else [50, 50, 50],
                        "acento": [0, 112, 192],
                        "destaque": [70, 130, 180]
                    },
                    "fontes": {"titulo": "Calibri", "corpo": "Calibri"},
                    "estilo_visual": "corporativo",
                    "elementos_especiais": ["blocos de cor"],
                    "layout": "estruturado",
                    "descricao": "Design azul corporativo profissional"
                }

        elif "escuro" in comando_lower or "dark" in comando_lower:
            return {
                "nome_estilo": "Dark Modern",
                "cores": {
                    "fundo": [30, 30, 35],
                    "titulo": [255, 255, 255],
                    "texto": [180, 180, 180],
                    "acento": [0, 255, 136],
                    "destaque": [100, 200, 255]
                },
                "fontes": {"titulo": "Calibri", "corpo": "Calibri"},
                "estilo_visual": "moderno",
                "elementos_especiais": ["neon sutil", "contrastes fortes"],
                "layout": "assimétrico",
                "descricao": "Design escuro moderno com acentos neon"
            }
        
        elif "minimalista" in comando_lower or "clean" in comando_lower:
            return {
                "nome_estilo": "Minimalista",
                "cores": {
                    "fundo": [255, 255, 255],
                    "titulo": [30, 30, 30],
                    "texto": [100, 100, 100],
                    "acento": [200, 200, 200],
                    "destaque": [50, 50, 50]
                },
                "fontes": {"titulo": "Helvetica", "corpo": "Helvetica"},
                "estilo_visual": "minimalista",
                "elementos_especiais": ["espaçamento generoso", "linhas finas"],
                "layout": "clean",
                "descricao": "Design minimalista clean e elegante"
            }
        
        # Detectar cor específica para header/footer
        cor_acento = [212, 175, 55]  # Dourado padrão
        if "verde" in comando_lower:
            cor_acento = [0, 128, 0]  # Verde
        elif "vermelho" in comando_lower:
            cor_acento = [220, 53, 69]  # Vermelho
        elif "azul" in comando_lower:
            cor_acento = [0, 112, 192]  # Azul
        elif "roxo" in comando_lower or "purple" in comando_lower:
            cor_acento = [128, 0, 128]  # Roxo
        elif "laranja" in comando_lower:
            cor_acento = [255, 165, 0]  # Laranja
        
        # Detectar fundo escuro/preto
        fundo_cor = [255, 255, 255]  # Branco padrão
        texto_cor = [45, 45, 45]     # Cinza escuro para texto

        # Se usuário quer texto preto explicitamente
        if usar_preto:
            fundo_cor = [255, 255, 255]  # Fundo branco
            texto_cor = [0, 0, 0]         # Texto preto puro
        elif "escuro" in comando_lower or "dark" in comando_lower:
            fundo_cor = [30, 30, 30]   # Quase preto
            texto_cor = [255, 255, 255]  # Texto branco
        elif "cinza" in comando_lower or "gray" in comando_lower or "grey" in comando_lower:
            fundo_cor = [60, 60, 60]   # Cinza escuro
            texto_cor = [255, 255, 255]  # Texto branco

        # Padrão: dourado e branco (Auditar)
        return {
            "nome_estilo": "Auditar Corporativo",
            "cores": {
                "fundo": fundo_cor,
                "titulo": texto_cor,
                "texto": [0, 0, 0] if usar_preto else ([200, 200, 200] if fundo_cor != [255, 255, 255] else [80, 80, 80]),
                "acento": cor_acento,
                "destaque": cor_acento
            },
            "fontes": {"titulo": "Calibri", "corpo": "Calibri"},
            "estilo_visual": "corporativo",
            "elementos_especiais": ["header_footer_colorido"],
            "layout": "profissional",
            "descricao": f"Design corporativo Auditar com fundo {'escuro' if fundo_cor != [255, 255, 255] else 'branco'} e acentos coloridos"
        }
    
    def _codigo_fallback(self, design, nome_empresa, responsavel, resumo, dados_mensais=None) -> str:
        """Gera código estático baseado no design escolhido"""
        
        cores = design['cores']
        fontes = design['fontes']
        
        # Preparar dados reais se disponíveis
        if dados_mensais and len(dados_mensais) > 0:
            primeiro = dados_mensais[0]
            ultimo = dados_mensais[-1]
            
            receita_total = sum(d.get("receita_bruta", 0) for d in dados_mensais)
            lucro_total = sum(d.get("lucro_operacional", 0) for d in dados_mensais)
            margem = (lucro_total / receita_total * 100) if receita_total > 0 else 0
            
            # Formatar mês/ano com tratamento de tipo
            def fmt_mes_ano(mes_val, ano_val, mes_default='01', ano_default='2024'):
                try:
                    mes = int(mes_val)
                    mes_str = f"{mes:02d}"
                except:
                    mes_str = str(mes_val).zfill(2)[:2] if mes_val else mes_default
                try:
                    ano = int(ano_val)
                    ano_str = str(ano)
                except:
                    ano_str = str(ano_val) if ano_val else ano_default
                return f"{mes_str}/{ano_str}"
            
            meses_str = fmt_mes_ano(primeiro.get('mes'), primeiro.get('ano'))
            if len(dados_mensais) > 1:
                meses_str += f" a {fmt_mes_ano(ultimo.get('mes'), ultimo.get('ano'))}"
            
            dados_resumo = {
                "primeiro_mes": meses_str.split(" a ")[0] if " a " in meses_str else meses_str,
                "ultimo_mes": meses_str.split(" a ")[1] if " a " in meses_str else meses_str,
                "receita_total": f"R$ {receita_total:,.2f}",
                "lucro_total": f"R$ {lucro_total:,.2f}",
                "margem": f"{margem:.1f}%"
            }
        else:
            # Fallback para dados de exemplo
            dados_resumo = {
                "primeiro_mes": "Jan/2024",
                "ultimo_mes": "Dez/2024",
                "receita_total": "R$ 0,00",
                "lucro_total": "R$ 0,00",
                "margem": "0.0%"
            }
        
        # Preparar dados para tabela mensal
        linhas_tabela = []
        if dados_mensais:
            for d in dados_mensais[:10]:
                # Tratar mês/ano que podem ser int ou string
                mes_val = d.get('mes', 0)
                ano_val = d.get('ano', 0)
                try:
                    mes_str = f"{int(mes_val):02d}"
                except:
                    mes_str = str(mes_val).zfill(2)[:2]
                try:
                    ano_str = str(int(ano_val))
                except:
                    ano_str = str(ano_val)
                mes = f"{mes_str}/{ano_str}"
                
                rec = f"R$ {d.get('receita_bruta', 0):,.0f}"
                custo_total = d.get('custos', 0) + d.get('despesas', 0) + d.get('impostos', 0)
                custo = f"R$ {custo_total:,.0f}"
                lucro = f"R$ {d.get('lucro_operacional', 0):,.0f}"
                linhas_tabela.append(f'{{"mes": "{mes}", "receita": "{rec}", "custo": "{custo}", "lucro": "{lucro}"}}')
        dados_tabela = ', '.join(linhas_tabela) if linhas_tabela else '{"mes": "N/A", "receita": "R$ 0", "custo": "R$ 0", "lucro": "R$ 0"}'
        
        return f'''"""
Gerador de Apresentação - Design: {design['nome_estilo']}
Código gerado automaticamente via IA Local
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Cores do design
COR_FUNDO = RGBColor({cores['fundo'][0]}, {cores['fundo'][1]}, {cores['fundo'][2]})
COR_TITULO = RGBColor({cores['titulo'][0]}, {cores['titulo'][1]}, {cores['titulo'][2]})
COR_TEXTO = RGBColor({cores['texto'][0]}, {cores['texto'][1]}, {cores['texto'][2]})
COR_ACENTO = RGBColor({cores['acento'][0]}, {cores['acento'][1]}, {cores['acento'][2]})
COR_DESTAQUE = RGBColor({cores['destaque'][0]}, {cores['destaque'][1]}, {cores['destaque'][2]})

FONTE_TITULO = "{fontes['titulo']}"
FONTE_CORPO = "{fontes['corpo']}"

def formatar_moeda(valor):
    return f"R$ {{valor:,.2f}}".replace(",", "X").replace(".", ",").replace("X", ".")

def adicionar_fundo(slide, cor=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = cor or COR_FUNDO

def adicionar_texto(slide, texto, x, y, w, h, tamanho=14, cor=None, bold=False, alinhamento=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alinhamento
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.color.rgb = cor or COR_TEXTO
    run.font.bold = bold
    run.font.name = FONTE_CORPO if not bold else FONTE_TITULO
    return txBox

def adicionar_logo(slide, x=0.3, y=4.85, largura=0.8):
    """Adiciona logo da Auditar ou texto fallback"""
    try:
        logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo_auditar.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(x), Inches(y), width=Inches(largura))
        else:
            # Fallback texto
            txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(largura), Inches(0.3))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "AUDITAR"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COR_ACENTO
    except:
        pass  # Ignora erro se não conseguir adicionar

def slide_capa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide)
    
    # Elemento decorativo
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()
    
    # Título
    adicionar_texto(slide, "PLANEJAMENTO TRIBUTÁRIO", 0.5, 1.0, 9, 1, 
                    tamanho=36, cor=COR_TITULO, bold=True, alinhamento=PP_ALIGN.CENTER)
    adicionar_texto(slide, "Análise Estratégica e Projeções", 0.5, 2.0, 9, 0.5,
                    tamanho=20, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)
    
    # Empresa
    adicionar_texto(slide, "{nome_empresa.upper()}", 0.5, 2.8, 9, 0.5,
                    tamanho=24, cor=COR_ACENTO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Info
    adicionar_texto(slide, "Responsável: {responsavel}", 0.5, 3.5, 9, 0.3,
                    tamanho=12, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)
    adicionar_texto(slide, "Emitido em: " + datetime.now().strftime("%d/%m/%Y"), 0.5, 3.8, 9, 0.3,
                    tamanho=12, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)
    
    # Logo
    adicionar_logo(slide, 0.3, 4.6, 1.0)
    
    # Footer
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.4), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()
    
    adicionar_texto(slide, "AUDITAR - CONTABILIDADE CONSULTIVA", 0.5, 5.0, 9, 0.3,
                    tamanho=11, cor=COR_ACENTO, bold=True, alinhamento=PP_ALIGN.CENTER)

def slide_resumo(prs, resumo_dados):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide)

    # Barra superior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

    # Título
    adicionar_texto(slide, "RESUMO EXECUTIVO", 0.3, 0.3, 9.4, 0.5,
                    tamanho=24, cor=COR_TITULO, bold=True)

    # Descrição
    adicionar_texto(slide, "Análise financeira consolidada do período",
                    0.3, 0.9, 9.4, 0.3, tamanho=14, cor=COR_TEXTO)

    # Cards com dados (sem barras coloridas)
    adicionar_texto(slide, "Período Analisado", 0.5, 1.5, 4, 0.3,
                    tamanho=12, cor=COR_TITULO, bold=True)
    adicionar_texto(slide, f"{{resumo_dados['primeiro_mes']}} a {{resumo_dados['ultimo_mes']}}",
                    0.5, 1.8, 4, 0.3, tamanho=14, cor=COR_TEXTO)

    adicionar_texto(slide, "Receita Total", 0.5, 2.3, 4, 0.3,
                    tamanho=12, cor=COR_TITULO, bold=True)
    adicionar_texto(slide, f"{{resumo_dados['receita_total']}}",
                    0.5, 2.6, 4, 0.4, tamanho=18, cor=COR_DESTAQUE, bold=True)

    adicionar_texto(slide, "Lucro Total", 5.5, 1.5, 4, 0.3,
                    tamanho=12, cor=COR_TITULO, bold=True)
    adicionar_texto(slide, f"{{resumo_dados['lucro_total']}}",
                    5.5, 1.8, 4, 0.4, tamanho=18, cor=COR_DESTAQUE, bold=True)

    adicionar_texto(slide, "Margem Média", 5.5, 2.3, 4, 0.3,
                    tamanho=12, cor=COR_TITULO, bold=True)
    adicionar_texto(slide, f"{{resumo_dados['margem']}}",
                    5.5, 2.6, 4, 0.4, tamanho=18, cor=COR_DESTAQUE, bold=True)

    # Barra inferior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.4), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

def slide_analise(prs, dados_empresa):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide)

    # Barra superior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

    # Título
    adicionar_texto(slide, "DADOS MENSAIS", 0.3, 0.3, 9.4, 0.5,
                    tamanho=24, cor=COR_TITULO, bold=True)

    # Cabeçalho da tabela
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.9), Inches(9.4), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

    adicionar_texto(slide, "Mês/Ano", 0.4, 0.95, 2, 0.3, tamanho=12, cor=COR_FUNDO, bold=True, alinhamento=PP_ALIGN.CENTER)
    adicionar_texto(slide, "Receita", 2.6, 0.95, 2.2, 0.3, tamanho=12, cor=COR_FUNDO, bold=True, alinhamento=PP_ALIGN.CENTER)
    adicionar_texto(slide, "Custos", 5.0, 0.95, 2, 0.3, tamanho=12, cor=COR_FUNDO, bold=True, alinhamento=PP_ALIGN.CENTER)
    adicionar_texto(slide, "Lucro", 7.2, 0.95, 2.2, 0.3, tamanho=12, cor=COR_FUNDO, bold=True, alinhamento=PP_ALIGN.CENTER)

    y_pos = 1.4
    for i, d in enumerate(dados_empresa):
        # Linha zebrada
        if i % 2 == 0:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(y_pos), Inches(9.4), Inches(0.35))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(248, 248, 248)
            shape.line.fill.background()

        adicionar_texto(slide, d['mes'], 0.4, y_pos + 0.05, 2, 0.25, tamanho=11, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)
        adicionar_texto(slide, d['receita'], 2.6, y_pos + 0.05, 2.2, 0.25, tamanho=11, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)
        adicionar_texto(slide, d['custo'], 5.0, y_pos + 0.05, 2, 0.25, tamanho=11, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)
        adicionar_texto(slide, d['lucro'], 7.2, y_pos + 0.05, 2.2, 0.25, tamanho=11, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)

        y_pos += 0.37

    # Barra inferior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.4), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

def slide_cenarios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide)

    # Barra superior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

    # Título
    adicionar_texto(slide, "CENÁRIOS PROJETADOS", 0.3, 0.3, 9.4, 0.5,
                    tamanho=24, cor=COR_TITULO, bold=True)

    # Cards de cenários (sem barras coloridas)
    cenarios = [
        ("OTIMISTA", "+30%", "Crescimento acelerado"),
        ("REALISTA", "+10%", "Crescimento estável"),
        ("PESSIMISTA", "-10%", "Cenário conservador")
    ]

    for i, (nome, perc, desc) in enumerate(cenarios):
        x = 0.3 + i * 3.2
        # Card sem barra colorida - apenas fundo e borda cinza
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.2), Inches(3.0), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COR_FUNDO if i == 1 else RGBColor(245, 245, 245)
        shape.line.color.rgb = RGBColor(200, 200, 200)

        adicionar_texto(slide, nome, x + 0.1, 1.25, 2.8, 0.3,
                        tamanho=12, cor=COR_TITULO, bold=True, alinhamento=PP_ALIGN.CENTER)
        adicionar_texto(slide, perc, x + 0.1, 1.9, 2.8, 0.4,
                        tamanho=20, cor=COR_DESTAQUE, bold=True, alinhamento=PP_ALIGN.CENTER)
        adicionar_texto(slide, desc, x + 0.1, 2.5, 2.8, 0.5,
                        tamanho=11, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)

    # Barra inferior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.4), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

def slide_recomendacoes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide)
    
    # Barra superior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()
    
    # Título
    adicionar_texto(slide, "RECOMENDAÇÕES ESTRATÉGICAS", 0.3, 0.3, 9.4, 0.5, 
                    tamanho=24, cor=COR_TITULO, bold=True)
    
    # Lista de recomendações
    recs = [
        "1. Revisão estrutural da carga tributária",
        "2. Análise de alíquotas por regime",
        "3. Otimização de custos operacionais",
        "4. Planejamento de fluxo de caixa",
        "5. Consultoria contínua mensal"
    ]
    
    y_pos = 1.2
    for rec in recs:
        adicionar_texto(slide, rec, 0.5, y_pos, 9, 0.4, 
                        tamanho=14, cor=COR_TEXTO)
        y_pos += 0.5
    
    # Barra inferior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.4), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

def slide_conclusao(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_fundo(slide)
    
    # Barra superior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()
    
    # Conteúdo centralizado
    adicionar_texto(slide, "AGRADECIMENTOS", 0.3, 1.5, 9.4, 0.8, 
                    tamanho=36, cor=COR_ACENTO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    adicionar_texto(slide, "Obrigado pela confiança", 0.3, 2.4, 9.4, 0.5, 
                    tamanho=20, cor=COR_TITULO, alinhamento=PP_ALIGN.CENTER)
    
    adicionar_texto(slide, "{nome_empresa.upper()}", 0.3, 3.0, 9.4, 0.5, 
                    tamanho=24, cor=COR_DESTAQUE, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    adicionar_texto(slide, "Responsável: {responsavel}", 0.3, 3.7, 9.4, 0.3, 
                    tamanho=12, cor=COR_TEXTO, alinhamento=PP_ALIGN.CENTER)
    
    adicionar_texto(slide, "AUDITAR - CONTABILIDADE CONSULTIVA", 0.3, 4.8, 9.4, 0.3, 
                    tamanho=11, cor=COR_ACENTO, bold=True, alinhamento=PP_ALIGN.CENTER)
    
    # Barra inferior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.4), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_ACENTO
    shape.line.fill.background()

def gerar_apresentacao():
    """Função principal - Gera apresentação completa de 6 slides"""
    # Estrutura:
    # 1. Capa
    # 2. Resumo Executivo (cards sem barras)
    # 3. Dados Mensais (tabela)
    # 4. Cenários (cards sem barras)
    # 5. Recomendações
    # 6. Conclusão

    # Preparar dados para os slides (valores reais da empresa)
    resumo_dados = dict(
        primeiro_mes="{{ dados_resumo['primeiro_mes'] }}",
        ultimo_mes="{{ dados_resumo['ultimo_mes'] }}",
        receita_total="{{ dados_resumo['receita_total'] }}",
        lucro_total="{{ dados_resumo['lucro_total'] }}",
        margem="{{ dados_resumo['margem'] }}"
    )

    # Dados da tabela mensal
    dados_empresa = [{dados_tabela}]

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # SLIDE 1: CAPA
    slide_capa(prs)

    # SLIDE 2: RESUMO EXECUTIVO (cards sem barras coloridas)
    slide_resumo(prs, resumo_dados)

    # SLIDE 3: DADOS MENSAIS (tabela)
    slide_analise(prs, dados_empresa)

    # SLIDE 4: CENÁRIOS (cards sem barras coloridas)
    slide_cenarios(prs)

    # SLIDE 5: RECOMENDAÇÕES
    slide_recomendacoes(prs)

    # SLIDE 6: CONCLUSÃO
    slide_conclusao(prs)
    
    # Salvar
    nome_arquivo = "Apresentacao_IA_{nome_empresa.replace(' ', '_')}_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".pptx"
    caminho = os.path.join(os.path.expanduser("~"), "Desktop", nome_arquivo)
    
    # Garantir diretório
    os.makedirs(os.path.dirname(caminho), exist_ok=True)
    
    prs.save(caminho)
    print(f"[OK] Apresentao salva em: {{caminho}}")
    return caminho

if __name__ == "__main__":
    gerar_apresentacao()
'''
    
    def _encontrar_arquivo_gerado(self, nome_empresa: str) -> str:
        """Procura por arquivo gerado recentemente"""
        return self._encontrar_arquivo_recente(nome_empresa, tempo_limite_segundos=60) or \
               os.path.join(os.path.expanduser("~"), "Desktop", 
                   f"Apresentacao_{nome_empresa.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    
    def _encontrar_arquivo_recente(self, nome_empresa: str, tempo_limite_segundos: int = 60) -> str:
        """Procura por arquivo PPTX criado recentemente no Desktop"""
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        
        try:
            # Listar arquivos PPTX
            pptx_files = [
                f for f in os.listdir(desktop)
                if f.endswith('.pptx')
            ]
            
            if not pptx_files:
                return None
            
            # Verificar arquivos criados nos últimos X segundos
            agora = datetime.now().timestamp()
            arquivos_recentes = []
            
            for f in pptx_files:
                caminho = os.path.join(desktop, f)
                try:
                    tempo_criacao = os.path.getctime(caminho)
                    if (agora - tempo_criacao) < tempo_limite_segundos:
                        arquivos_recentes.append((caminho, tempo_criacao))
                except:
                    continue
            
            if arquivos_recentes:
                # Retornar o mais recente
                arquivos_recentes.sort(key=lambda x: x[1], reverse=True)
                return arquivos_recentes[0][0]
            
            return None
            
        except Exception as e:
            print(f"[AVISO]  Erro ao procurar arquivo: {e}")
            return None
    
    def _gerar_fallback_local(
        self,
        dados_mensais: List[Dict],
        nome_empresa: str,
        responsavel: str,
        design: Dict
    ) -> str:
        """Gera apresentação usando código local quando IA falha"""
        
        print("[DESIGN] Usando gerador local com design:", design['nome_estilo'])
        
        # Criar apresentação básica
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Extrair cores
        cores = design['cores']
        COR_FUNDO = RGBColor(*cores['fundo'])
        COR_TITULO = RGBColor(*cores['titulo'])
        COR_TEXTO = RGBColor(*cores['texto'])
        COR_ACENTO = RGBColor(*cores['acento'])
        
        # Funções auxiliares
        def add_fundo(slide, cor=COR_FUNDO):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = cor
        
        def add_texto(slide, texto, x, y, w, h, tamanho=14, cor=COR_TEXTO, bold=False, align=PP_ALIGN.LEFT):
            tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tx.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = texto
            run.font.size = Pt(tamanho)
            run.font.color.rgb = cor
            run.font.bold = bold
        
        def add_logo(slide, x=0.3, y=4.85, w=0.8):
            """Adiciona logo da Auditar ou texto"""
            try:
                import os
                logo_path = os.path.join(os.path.dirname(__file__), "logo_auditar.png")
                if os.path.exists(logo_path):
                    slide.shapes.add_picture(logo_path, Inches(x), Inches(y), width=Inches(w))
                else:
                    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
                    tf = tx.text_frame
                    p = tf.paragraphs[0]
                    p.text = "AUDITAR"
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = COR_ACENTO
            except:
                pass
        
        def add_retangulo(slide, x, y, w, h, cor):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = cor
            shape.line.fill.background()
        
        # SLIDE 1: CAPA
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_fundo(slide)
        add_retangulo(slide, 0, 0, 10, 0.15, COR_ACENTO)
        add_texto(slide, "PLANEJAMENTO TRIBUTÁRIO", 0.3, 1.0, 9.4, 0.8, 32, COR_TITULO, True, PP_ALIGN.CENTER)
        add_texto(slide, nome_empresa.upper(), 0.3, 2.0, 9.4, 0.5, 24, COR_TITULO, True, PP_ALIGN.CENTER)
        add_texto(slide, f"Responsável: {responsavel}", 0.3, 3.0, 9.4, 0.3, 12, COR_TEXTO, False, PP_ALIGN.CENTER)
        add_logo(slide, 0.3, 4.6, 1.0)
        add_retangulo(slide, 0, 5.4, 10, 0.15, COR_ACENTO)
        
        # Calcular indicadores
        receita_total = sum(d["receita_bruta"] for d in dados_mensais)
        lucro_total = sum(d["lucro_operacional"] for d in dados_mensais)
        margem = (lucro_total / receita_total * 100) if receita_total > 0 else 0

        # SLIDE 2: RESUMO EXECUTIVO
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_fundo(slide)
        add_retangulo(slide, 0, 0, 10, 0.15, COR_ACENTO)
        add_texto(slide, "RESUMO EXECUTIVO", 0.3, 0.3, 9.4, 0.5, 24, COR_TITULO, True)

        # Cards de indicadores (sem barra dourada)
        indicadores = [
            ("Receita Total", f"R$ {receita_total:,.2f}"),
            ("Lucro Total", f"R$ {lucro_total:,.2f}"),
            ("Margem", f"{margem:.1f}%")
        ]

        for i, (titulo, valor) in enumerate(indicadores):
            x = 0.5 + i * 3.2
            add_retangulo(slide, x, 1.5, 2.8, 1.5, COR_FUNDO)
            add_texto(slide, titulo, x + 0.1, 1.6, 2.6, 0.3, 12, COR_TITULO, True, PP_ALIGN.CENTER)
            add_texto(slide, valor, x + 0.1, 2.0, 2.6, 0.5, 16, COR_TITULO, True, PP_ALIGN.CENTER)

        add_retangulo(slide, 0, 5.4, 10, 0.15, COR_ACENTO)

        # SLIDE 3: DADOS MENSAIS DETALHADOS
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_fundo(slide)
        add_retangulo(slide, 0, 0, 10, 0.15, COR_ACENTO)
        add_texto(slide, "DADOS MENSAIS DETALHADOS", 0.3, 0.3, 9.4, 0.5, 22, COR_TITULO, True)

        # Cabeçalho da tabela
        add_retangulo(slide, 0.3, 0.9, 9.4, 0.4, COR_ACENTO)
        add_texto(slide, "Mês/Ano", 0.4, 0.95, 2, 0.3, 12, COR_FUNDO, True, PP_ALIGN.CENTER)
        add_texto(slide, "Receita", 2.6, 0.95, 2.2, 0.3, 12, COR_FUNDO, True, PP_ALIGN.CENTER)
        add_texto(slide, "Custos", 5.0, 0.95, 2, 0.3, 12, COR_FUNDO, True, PP_ALIGN.CENTER)
        add_texto(slide, "Lucro", 7.2, 0.95, 2.2, 0.3, 12, COR_FUNDO, True, PP_ALIGN.CENTER)

        # Dados mensais (mostra até 10 meses)
        y_pos = 1.4
        for i, d in enumerate(dados_mensais[:10]):
            # Tratar mês/ano que podem ser int ou string
            mes_val = d.get('mes', 0)
            ano_val = d.get('ano', 0)
            try:
                mes_str = f"{int(mes_val):02d}"
            except:
                mes_str = str(mes_val).zfill(2)[:2]
            try:
                ano_str = str(int(ano_val))
            except:
                ano_str = str(ano_val)
            mes = f"{mes_str}/{ano_str}"

            rec = f"R$ {d.get('receita_bruta', 0):,.0f}"
            custo = f"R$ {d.get('custos', 0) + d.get('despesas', 0) + d.get('impostos', 0):,.0f}"
            lucro = f"R$ {d.get('lucro_operacional', 0):,.0f}"

            # Linha zebrada
            if i % 2 == 0:
                add_retangulo(slide, 0.3, y_pos, 9.4, 0.35, RGBColor(248, 248, 248))

            add_texto(slide, mes, 0.4, y_pos + 0.05, 2, 0.25, 11, COR_TEXTO, False, PP_ALIGN.CENTER)
            add_texto(slide, rec, 2.6, y_pos + 0.05, 2.2, 0.25, 11, COR_TEXTO, False, PP_ALIGN.CENTER)
            add_texto(slide, custo, 5.0, y_pos + 0.05, 2, 0.25, 11, COR_TEXTO, False, PP_ALIGN.CENTER)
            add_texto(slide, lucro, 7.2, y_pos + 0.05, 2.2, 0.25, 11, COR_TEXTO, False, PP_ALIGN.CENTER)

            y_pos += 0.37

        add_retangulo(slide, 0, 5.4, 10, 0.15, COR_ACENTO)

        # SLIDE 4: CENÁRIOS
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_fundo(slide)
        add_retangulo(slide, 0, 0, 10, 0.15, COR_ACENTO)
        add_texto(slide, "ANÁLISE DE CENÁRIOS", 0.3, 0.3, 9.4, 0.5, 24, COR_TITULO, True)

        cenarios = [
            ("OTIMISTA", "+30%", receita_total * 1.3, lucro_total * 1.3),
            ("REALISTA", "+10%", receita_total * 1.1, lucro_total * 1.1),
            ("PESSIMISTA", "-20%", receita_total * 0.8, lucro_total * 0.8)
        ]

        for i, (nome, var, rec, luc) in enumerate(cenarios):
            x = 0.3 + i * 3.2
            add_retangulo(slide, x, 1.2, 3.0, 2.5, COR_FUNDO)
            add_texto(slide, nome, x + 0.1, 1.25, 2.8, 0.3, 12, COR_TITULO, True, PP_ALIGN.CENTER)
            add_texto(slide, f"Variação: {var}", x + 0.1, 1.7, 2.8, 0.3, 11, COR_TEXTO)
            add_texto(slide, f"R$ {rec:,.0f}", x + 0.1, 2.1, 2.8, 0.3, 14, COR_TITULO, True)

        add_retangulo(slide, 0, 5.4, 10, 0.15, COR_ACENTO)

        # SLIDE 5: RECOMENDAÇÕES
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_fundo(slide)
        add_retangulo(slide, 0, 0, 10, 0.15, COR_ACENTO)
        add_texto(slide, "RECOMENDAÇÕES ESTRATÉGICAS", 0.3, 0.3, 9.4, 0.5, 22, COR_TITULO, True)

        recomendacoes = [
            "1. Revisão estrutural da carga tributária atual",
            "2. Análise comparativa de regimes tributários",
            "3. Otimização de custos operacionais dedutíveis",
            "4. Planejamento de fluxo de caixa tributário",
            "5. Consultoria contínua para compliance fiscal"
        ]

        y = 1.2
        for rec in recomendacoes:
            add_texto(slide, rec, 0.5, y, 9, 0.4, 13, COR_TEXTO)
            y += 0.55

        add_retangulo(slide, 0, 5.4, 10, 0.15, COR_ACENTO)

        # SLIDE 6: CONCLUSÃO
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_fundo(slide)
        add_retangulo(slide, 0, 0, 10, 0.15, COR_ACENTO)
        add_texto(slide, "AGRADECIMENTOS", 0.3, 1.5, 9.4, 0.8, 36, COR_TITULO, True, PP_ALIGN.CENTER)
        add_texto(slide, nome_empresa.upper(), 0.3, 2.5, 9.4, 0.5, 24, COR_TITULO, True, PP_ALIGN.CENTER)
        add_texto(slide, f"Responsável: {responsavel}", 0.3, 3.3, 9.4, 0.3, 12, COR_TEXTO, False, PP_ALIGN.CENTER)
        add_logo(slide, 0.3, 4.5, 1.0)
        add_texto(slide, "AUDITAR - CONTABILIDADE CONSULTIVA", 0.3, 4.8, 9.4, 0.3, 11, COR_TITULO, True, PP_ALIGN.CENTER)
        add_retangulo(slide, 0, 5.4, 10, 0.15, COR_ACENTO)
        
        # Salvar
        nome_arquivo = f"Apresentacao_IA_{nome_empresa.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        caminho = os.path.join(os.path.expanduser("~"), "Desktop", nome_arquivo)
        os.makedirs(os.path.dirname(caminho), exist_ok=True)
        prs.save(caminho)
        
        return caminho


# Importar templates
from slide_templates import TEMPLATES, aplicar_template


# Função de compatibilidade com o sistema existente - USA TEMPLATES PRÉ-FABRICADOS
async def gerar_apresentacao_ia(
    dados_mensais: List[Dict],
    nome_empresa: str,
    responsavel: str,
    comando_estilo: str = "",
    model: str = DEFAULT_MODEL,
    bundle_dir: str = None,
    cores_personalizadas: dict = None
) -> str:
    """
    Função principal para integração com o sistema contábil
    USA TEMPLATES PRÉ-FABRICADOS - mais robusto que gerar código

    Args:
        dados_mensais: Lista de dicionários com dados financeiros
        nome_empresa: Nome da empresa
        responsavel: Nome do responsável
        comando_estilo: Comando em linguagem natural (ex: "fundo azul marinho com dourado")
        model: Modelo Ollama a usar
        cores_personalizadas: Dicionário de cores personalizadas (opcional)

    Returns:
        Caminho do arquivo PPTX gerado
    """
    print(f"\nAnalisando estilo: '{comando_estilo}'")
    print(f"Dados: {len(dados_mensais)} meses de {nome_empresa}")

    # IA apenas escolhe o template e as cores (não gera código)
    gerador = GeradorIAInteligente(model=model)

    if not comando_estilo:
        comando_estilo = "corporativo escuro azul e dourado"

    # Se cores personalizadas foram fornecidas, usar elas
    if cores_personalizadas:
        print("Usando cores personalizadas fornecidas pelo usurio")
        # Ainda usar IA para escolher o template, mas não para cores
        try:
            design = await gerador._analisar_requisitos(comando_estilo)
        except Exception as e:
            print(f"Erro ao analisar requisitos: {e}")
            design = gerador._design_fallback(comando_estilo)
    else:
        # Analisar comando para escolher template e cores
        try:
            design = await gerador._analisar_requisitos(comando_estilo)
        except Exception as e:
            print(f"Erro ao analisar requisitos: {e}")
            # Usar design fallback em caso de erro
            design = gerador._design_fallback(comando_estilo)

    # Mapear design para template (Claro ou Escuro)
    cmd_l = comando_estilo.lower()
    is_claro = design["cores"]["fundo"] == [255, 255, 255] or "branco" in cmd_l or "claro" in cmd_l
    is_escuro = "escuro" in cmd_l or "dark" in cmd_l or "preto" in cmd_l
    
    if is_escuro:
        template_nome = "corporativo_escuro"
        fundo_rgb = [33, 37, 41] # Cinza bem escuro profissional
    elif is_claro:
        template_nome = "corporativo_claro"
        fundo_rgb = [255, 255, 255]
    else:
        template_nome = "corporativo_escuro" # Padrao
        fundo_rgb = [33, 37, 41]

    # Se cores personalizadas não foram fornecidas, gerar com IA
    if not cores_personalizadas:
        # Extrair cores do design retornado (fallback ou IA)
        cores_design = design.get("cores", {})
        acento_ia = cores_design.get("acento", [212, 175, 55])
        
        # Dicionário de cores para segurança (suporta quase todas as cores comuns)
        mapa_cores = {
            "vermelho": [200, 0, 0], "azul": [0, 80, 180], "verde": [34, 139, 34],
            "amarelo": [255, 215, 0], "roxo": [128, 0, 128], "rosa": [255, 105, 180],
            "laranja": [255, 140, 0], "cinza": [128, 128, 128], "preto": [0, 0, 0],
            "branco": [255, 255, 255], "ciano": [0, 255, 255], "magenta": [255, 0, 255],
            "marrom": [139, 69, 19], "dourado": [212, 175, 55], "prata": [192, 192, 192],
            "flamengo": [200, 0, 0], "palmeiras": [0, 100, 0], "corinthians": [30, 30, 30],
            "gremio": [0, 150, 255], "inter": [200, 0, 0], "cruzeiro": [0, 0, 200]
        }
        
        cor_destaque = acento_ia
        # Verificar se alguma cor do dicionário está no comando do usuário
        for nome_cor, rgb in mapa_cores.items():
            if nome_cor in cmd_l:
                cor_destaque = rgb
                break

        # Criar cores personalizadas
        from pptx.dml.color import RGBColor
        cores_personalizadas = {
            'primaria': RGBColor(*cor_destaque), 
            'secundaria': RGBColor(*cor_destaque), 
            'fundo': RGBColor(*fundo_rgb),       
            'texto': RGBColor(255, 255, 255) if is_escuro else RGBColor(33, 37, 41),
            'texto_secundario': RGBColor(200, 200, 200) if is_escuro else RGBColor(108, 117, 125),
            'destaque': RGBColor(*cor_destaque),
            'accent': RGBColor(*cor_destaque)
        }

        print(f"[DESIGN] Cores aplicadas: acento RGB({', '.join(map(str, cor_destaque))})")
    else:
        print("[DESIGN] Cores personalizadas aplicadas")

    # Escolher template baseado nas preferências
    estilo = design.get("estilo", "").lower()
    cores_escolhidas = design.get("cores", {})

    if "minimalista" in estilo or "branco" in estilo or "clean" in estilo:
        template_nome = "minimalista_branco"
        # Cores padrão do minimalista podem ser sobrescritas
        if cores_escolhidas:
            from pptx.dml.color import RGBColor
            cores_personalizadas = {
                'primaria': RGBColor(51, 51, 51),
                'secundaria': RGBColor(220, 53, 69),  # Vermelho padrão
                'fundo': RGBColor(255, 255, 255),
                'texto': RGBColor(33, 37, 41),
                'texto_secundario': RGBColor(108, 117, 125),
                'destaque': RGBColor(40, 167, 69)  # Verde
            }
            # Personalizar se IA especificou cores
            cor_secundaria = cores_escolhidas.get("secundaria", "")
            if "vermelho" in cor_secundaria.lower():
                cores_personalizadas['secundaria'] = RGBColor(220, 53, 69)
            elif "verde" in cor_secundaria.lower():
                cores_personalizadas['secundaria'] = RGBColor(40, 167, 69)
                cores_personalizadas['destaque'] = RGBColor(220, 53, 69)
            elif "azul" in cor_secundaria.lower():
                cores_personalizadas['secundaria'] = RGBColor(0, 123, 255)
                cores_personalizadas['destaque'] = RGBColor(40, 167, 69)

    elif "moderno" in estilo or "gradiente" in estilo or "vibrante" in estilo:
        template_nome = "moderno_gradiente"

    # Se escolheu corporativo escuro mas pediu cores específicas
    if template_nome == "corporativo_escuro" and cores_escolhidas:
        from pptx.dml.color import RGBColor
        primaria = cores_escolhidas.get("primaria", "").lower()

        if "preto" in primaria or "escuro" in primaria or "cinza" in primaria:
            cores_personalizadas = {
                'primaria': RGBColor(30, 30, 30),
                'secundaria': RGBColor(220, 53, 69) if "vermelho" in str(cores_escolhidas) else RGBColor(40, 167, 69),
                'fundo': RGBColor(20, 20, 20),
                'texto': RGBColor(255, 255, 255),
                'texto_secundario': RGBColor(200, 200, 200),
                'destaque': RGBColor(220, 53, 69) if "vermelho" in str(cores_escolhidas) else RGBColor(212, 175, 55)
            }

    print(f"[DESIGN] Template escolhido: {template_nome}")

    # Criar apresentação usando template
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    aplicar_template(template_nome, prs, dados_mensais, nome_empresa, responsavel, cores_personalizadas)

    # Salvar
    nome_arquivo = f"Apresentacao_IA_{nome_empresa.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    caminho = os.path.join(os.path.expanduser("~"), "Desktop", nome_arquivo)
    os.makedirs(os.path.dirname(caminho), exist_ok=True)
    prs.save(caminho)

    print(f"[OK] Apresentao gerada: {caminho}")
    return caminho


# Execução standalone
if __name__ == "__main__":
    import asyncio
    
    # Teste
    dados_teste = [
        {"mes": 1, "ano": 2024, "receita_bruta": 100000, "custos": 60000, 
         "despesas": 20000, "impostos": 5000, "lucro_operacional": 15000}
    ]
    
    async def testar():
        resultado = await gerar_apresentacao_ia(
            dados_mensais=dados_teste,
            nome_empresa="Empresa Teste",
            responsavel="Contador",
            comando_estilo="fundo azul marinho escuro com dourado elegante, estilo luxuoso"
        )
        print(f"\n Resultado: {resultado}")
    
    asyncio.run(testar())
